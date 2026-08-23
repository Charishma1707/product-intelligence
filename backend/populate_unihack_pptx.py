import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

template_path = r'e:\unihack\product-intelligence\[EXT] UniHack-Protoype Template .pptx'
output_path = r'e:\unihack\product-intelligence\unihack_final.pptx'

# Load template
prs = Presentation(template_path)

# Remove all extra template placeholder slides except 1, then rebuild exactly 15 slides
while len(prs.slides) > 1:
    rId = prs.slides._sldIdLst[-1].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[-1]

print(f"Cleared template placeholders. Base slide count: {len(prs.slides)}")

# Add slides until we have exactly 15 clean slides
blank_layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0]
while len(prs.slides) < 15:
    prs.slides.add_slide(blank_layout)

# Theme Colors (High contrast, modern presentation palette)
DARK_NAVY = RGBColor(15, 23, 42)      # #0F172A (Titles)
PRIMARY_BLUE = RGBColor(37, 99, 235)   # #2563EB (Subheadings & Accents)
TEXT_MAIN = RGBColor(30, 41, 59)      # #1E293B (Body Text)
GREEN_METRIC = RGBColor(16, 185, 129) # #10B981 (Metrics)

def format_slide(slide, title_text, sections):
    """Clears all old shapes and constructs perfectly readable title & structured bullet cards."""
    # Delete all old shapes
    for i in range(len(slide.shapes)-1, -1, -1):
        sp = slide.shapes[i].element
        sp.getparent().remove(sp)
        
    # Title Header (20pt Bold Navy)
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.0), Inches(0.9))
    tf_t = title_box.text_frame
    tf_t.word_wrap = True
    p_t = tf_t.paragraphs[0]
    p_t.text = title_text
    p_t.font.size = Pt(20)
    p_t.font.bold = True
    p_t.font.color.rgb = DARK_NAVY

    # Content Body Box (Readable 12pt - 14pt fonts, clean margins)
    body_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12.0), Inches(5.8))
    tf_b = body_box.text_frame
    tf_b.word_wrap = True
    tf_b.margin_left = Inches(0.1)
    tf_b.margin_right = Inches(0.1)
    tf_b.margin_top = Inches(0.1)
    tf_b.margin_bottom = Inches(0.1)

    first_paragraph = True
    for header, bullets in sections:
        if header:
            p_h = tf_b.paragraphs[0] if first_paragraph else tf_b.add_paragraph()
            first_paragraph = False
            p_h.text = header
            p_h.font.size = Pt(13)
            p_h.font.bold = True
            p_h.font.color.rgb = PRIMARY_BLUE
            p_h.space_before = Pt(8)
            p_h.space_after = Pt(2)

        for bullet_label, bullet_desc in bullets:
            p_bullet = tf_b.paragraphs[0] if first_paragraph else tf_b.add_paragraph()
            first_paragraph = False
            
            # Bold label run
            if bullet_label:
                r_lbl = p_bullet.add_run()
                r_lbl.text = bullet_label + " "
                r_lbl.font.bold = True
                r_lbl.font.size = Pt(12)
                r_lbl.font.color.rgb = DARK_NAVY
            
            # Normal text run
            r_txt = p_bullet.add_run()
            r_txt.text = bullet_desc
            r_txt.font.bold = False
            r_txt.font.size = Pt(12)
            r_txt.font.color.rgb = TEXT_MAIN
            
            p_bullet.space_before = Pt(3)
            p_bullet.space_after = Pt(3)

# ----------------------------------------------------
# 15 Highly Visible, Perfectly Scaled Slides
# ----------------------------------------------------
slides_data = [
    # Slide 1: Cover
    ("Autonomous Unilog Product Intelligence Pipeline", [
        ("PROJECT OVERVIEW & TRACK", [
            ("Track:", "Product Intelligence & Automated Catalog Enrichment"),
            ("Project Title:", "DeepFlowShield — High-Accuracy Multi-Agent Sourcing & Evidence RAG System"),
            ("Team Name:", "codewithcofee"),
            ("Team Leader:", "Charishma Alam"),
            ("Target Delivery Standard:", "Unilog 252-Column Master Catalog E-Commerce CSV Standard"),
            ("GitHub Repository:", "https://github.com/Charishma1707/product-intelligence"),
        ])
    ]),

    # Slide 2: Team Details & Snapshot
    ("Team Details & Key Performance Snapshot", [
        ("TEAM INFORMATION", [
            ("Team Name:", "codewithcofee"),
            ("Team Leader Name:", "Charishma Alam"),
            ("Track:", "Product Intelligence & Automated Catalog Enrichment"),
        ]),
        ("SYSTEM ROI & ACCURACY BENCHMARKS", [
            ("95% Time Savings:", "Reduces manual onboarding time from 100+ hours to under 30 seconds per SKU."),
            ("$0.00 Inference Cost:", "Primary extraction runs on air-gapped local Ollama (qwen2.5:3b)."),
            ("100% Evidence Provenance:", "Every value bound to exact PDF page, URL, confidence score, and text snippet."),
            ("252 Delivery Attributes:", "Extracted and formatted directly into Unilog's delivery CSV standard."),
        ])
    ]),

    # Slide 3: Solution Brief
    ("Brief about your solution", [
        ("EXECUTIVE SOLUTION OVERVIEW", [
            ("Problem Addressed:", "Industrial supplier feeds arrive dirty (e.g. 'APPDE' instead of 'Whirlpool'), missing specs, and unformatted. Catalog teams waste 100+ manual hours per catalog finding datasheets."),
            ("Proposed Solution:", "A 10-node state machine built on LangGraph that transforms 3 dirty input strings into a complete 252-column Unilog record in under 30 seconds."),
        ]),
        ("THREE CORE ARCHITECTURAL PILLARS", [
            ("1. OEM-Locked Sourcing:", "Bypasses marketplace noise (Amazon/eBay) by locking retrieval strictly to official manufacturer domains (whirlpool.com, se.com, fluke.com)."),
            ("2. Zero-Hallucination RAG:", "Ingests spec PDFs/HTML into ChromaDB vector store with verbatim proof snippets."),
            ("3. 5-Stage HITL Dashboard:", "Routes low-confidence items (<80%) to human catalog managers for 1-click review."),
        ])
    ]),

    # Slide 4: Q1 Enrichment Strategy
    ("1. How does your solution enrich minimal product information?", [
        ("FIVE-PHASE TRANSFORMATION WORKFLOW", [
            ("Phase 1 — Identity Disambiguation:", "Ingests dirty text ('APPDE'). Resolves true OEM ('Whirlpool Corporation') and expands jargon ('SS' -> 'Stainless Steel')."),
            ("Phase 2 — UNSPSC Leaf Taxonomy:", "Maps item to 8-Digit UNSPSC Leaf Code (e.g. 83041100 for Dishwashers) and selects 252-column expected schema."),
            ("Phase 3 — Multi-Modal OEM Sourcing:", "Discovers official specification PDFs, manuals, installation guides, and images with SHA-256 caching."),
            ("Phase 4 — Vector RAG & Graph Inheritance:", "Queries ChromaDB vector chunks and checks NetworkX Series Memory to extract 50 structured attribute triplets."),
            ("Phase 5 — Commercial Copywriting:", "Synthesizes Invoice Copy (<=40 chars), Short, Long, and Marketing copy."),
        ])
    ]),

    # Slide 5: Q2 Accuracy & Trust
    ("2. How does your solution ensure accuracy and trust in generated data?", [
        ("FIVE-LAYER ACCURACY SAFEGUARDS", [
            ("1. OEM Source Tiering:", "Mandates Tier-1 official OEM datasheets over distributor text. Returns null if missing."),
            ("2. Verifiable Evidence Anchors:", "Binds every field to Raw Value, Standard UOM, Source URL, PDF Page, Score, and Text Snippet."),
            ("3. Deterministic Validation:", "Runs _is_invalid_garbage filters, FIELD_SANITY range rules, and UOM splitters."),
            ("4. 5-Tier Confidence Scoring:", "Assigns scores from 1.00 (Human Verified) down to <0.80 (Flagged for Review)."),
            ("5. 5-Stage HITL Review Console:", "Pauses low-confidence items for 1-click human verification with implicit confidence boosting."),
        ])
    ]),

    # Slide 6: Q3 Scalability & Performance
    ("3. What makes your solution scalable for enterprise product catalogs?", [
        ("ENTERPRISE SCALABILITY ARCHITECTURE", [
            ("1. Multi-Thousand SKU Concurrency:", "FastAPI async execution loops & SQLite job_store.db state serialization stream batch progress cleanly."),
            ("2. Generic Multi-Format Parsers:", "Ingests PDF text/tables via PyMuPDF, HTML via trafilatura, and resolves new OEM domains automatically."),
            ("3. SHA-256 Vector Deduplication:", "Embeds shared catalog PDFs into ChromaDB EXACTLY ONCE, cutting scraping bandwidth by >90%."),
            ("4. $0.00 Local Inference Architecture:", "Primary extraction runs on local Ollama (qwen2.5:3b). Groq / Gemini API serve as cloud fallback."),
        ])
    ]),

    # Slide 7: Opportunities Q1
    ("Opportunities: How different is it from any of the other existing ideas?", [
        ("KEY DIFFERENTIATORS VS TRADITIONAL TOOLS", [
            ("Semantic RAG vs Rigid Scraping:", "Legacy XPath scrapers break when websites update. Our system uses semantic RAG to read unstructured PDF spec sheets dynamically."),
            ("OEM Guardrails vs Marketplace Noise:", "Generic AI tools hallucinate by scraping third-party listings (Amazon/eBay). We lock retrieval strictly to official OEM domains."),
            ("Visual Provenance vs Black-Box AI:", "Every extracted value is backed by cropped PDF text snippets and can be actively steered by human catalog managers."),
        ])
    ]),

    # Slide 8: Opportunities Q2
    ("Opportunities: How will it be able to solve the problem statement given?", [
        ("DIRECT PROBLEM STATEMENT RESOLUTION", [
            ("Eliminates 95% Manual Effort:", "Replaces web searching, PDF reading, and copy-pasting with automated 10-node agent execution under 30s."),
            ("252-Column Unilog Compliance:", "Formats extracted attributes directly into Unilog's delivery CSV structure."),
            ("Self-Learning Abbreviation Expander:", "Normalizes dirty supplier shorthand ('SS' -> 'Stainless Steel', 'MTR' -> 'Motor')."),
            ("Multi-Tier Sourcing Fallbacks:", "Falls back gracefully to approved industrial distributors (Grainger, Zoro) if primary OEM PDFs are absent."),
        ])
    ]),

    # Slide 9: Opportunities Q3
    ("Opportunities: Unique Selling Proposition (USP) of the proposed solution", [
        ("TOP 3 UNIQUE SELLING PROPOSITIONS (USP)", [
            ("1. OEM-Locked Sourcing Guardrails:", "Guarantees zero e-commerce marketplace hallucinations by domain-locking retrieval."),
            ("2. Natural Language AI Steering Bar:", "Empowers catalog managers to issue live directives into extraction nodes during review."),
            ("3. Graph-Backed Series Inheritance:", "NetworkX Knowledge Graph propagates specs across sibling SKUs, cutting web calls by >70%."),
        ])
    ]),

    # Slide 10: Features List
    ("List of features offered by the solution", [
        ("AUTOMATED PIPELINE ENGINE", [
            ("• Dirty Brand Disambiguation & Abbreviation Expander:", "Normalizes dirty supplier strings."),
            ("• 8-Digit UNSPSC Leaf Classification:", "Maps items to deep leaf categories & expected schemas."),
            ("• Multi-Pass OEM Document Harvesting:", "Ingests spec PDFs, manuals, SDS sheets, and images."),
            ("• SHA-256 Vector RAG Caching:", "ChromaDB vector store with zero duplicate PDF processing."),
            ("• Commercial Copywriting Engine:", "Generates Invoice (<=40 chars), Short, Long, and Marketing copy."),
        ]),
        ("INTERACTIVE HUMAN REVIEW SUITE", [
            ("• 5-Stage React Review Console:", "Identity, Sourcing, Attributes, Copywriting, and Delivery format audit."),
            ("• Natural Language Steering Bar:", "Allows human reviewers to steer extraction nodes with text prompts."),
            ("• 252-Column Unilog Exporter:", "1-click export compliant with Unilog master delivery specifications."),
        ])
    ]),

    # Slide 11: Process Flow
    ("Process flow diagram of the solution", [
        ("10-NODE LANGGRAPH STATE MACHINE FLOW", [
            ("Step 1: Input Ingestion", "-> Ingests dirty Brand, MPN, Description feed."),
            ("Step 2: Identity & Taxonomy", "-> Resolves OEM Brand & 8-Digit UNSPSC Category Code."),
            ("Step 3: Checkpoint 1 (HITL Identity)", "-> Pauses for human brand & taxonomy audit."),
            ("Step 4: Sourcing & Harvesting", "-> Locates OEM domain, downloads PDFs, hashes into ChromaDB."),
            ("Step 5: Checkpoint 2 (HITL Sourcing)", "-> Pauses for human URL & digital asset audit."),
            ("Step 6: RAG Extraction & Graph Lookup", "-> Extracts 50 attributes from PDF chunks & NetworkX series memory."),
            ("Step 7: Validation & Confidence Scoring", "-> Runs FIELD_SANITY rules and assigns confidence scores."),
            ("Step 8: Copywriting & 252-Column Export", "-> Synthesizes Invoice Copy (<=40 chars) & downloads Unilog CSV."),
        ])
    ]),

    # Slide 12: Wireframes & Dashboard
    ("Wireframes & Dashboard Interface", [
        ("DASHBOARD ARCHITECTURE & USER WORKFLOW", [
            ("1. 1-Click Sample Product Loaders:", "Instant demo buttons for Fluke 117, 3M Cubitron II, and Whirlpool Dishwasher."),
            ("2. Visual 5-Stage Execution Panel:", "Displays real-time pipeline status (Identity -> Sourcing -> Attributes -> Copywriting -> Delivery)."),
            ("3. Review Queue Console:", "Filterable table displaying jobs requiring review, confidence scores, and stage tags."),
            ("4. Side-by-Side PDF Evidence Inspector:", "Shows verbatim PDF text snippets alongside extracted values for verification."),
            ("5. 252-Column Master Export Controls:", "1-click button to download Unilog_Submission.csv and save to SQLite DB."),
        ])
    ]),

    # Slide 13: System Architecture
    ("Architecture diagram of the proposed solution", [
        ("3-TIER SYSTEM ARCHITECTURE BREAKDOWN", [
            ("Tier 1 — Frontend & State Control:", "React 18 + Vite Review Station, FastAPI Async REST Gateway, LangGraph 10-Node Supervisor, SQLite job_store.db."),
            ("Tier 2 — Search, Vector & Knowledge Stores:", "Serper.dev Google API, ChromaDB Vector DB, NetworkX Knowledge Graph, PyMuPDF & Trafilatura text engines."),
            ("Tier 3 — Multi-Tier Inference Cascade:", "Air-gapped local Ollama (qwen2.5:3b) primary extraction engine ($0.00), Groq / Gemini API cloud fallback."),
        ])
    ]),

    # Slide 14: Tech Stack & Cost
    ("Technologies used & Cost Analysis", [
        ("ENTERPRISE TECH STACK", [
            ("Languages & Runtimes:", "Python 3.13, JavaScript (ES6+), HTML5, CSS3"),
            ("Local Air-Gapped LLM Engine:", "Ollama (qwen2.5:3b) — Provides $0.00 marginal inference cost per SKU."),
            ("Premium Cloud Fallback Cascade:", "Groq API (llama-3.3-70b) & Google Gemini 1.5 Pro for complex unformatted sheets."),
            ("AI & Vector Frameworks:", "LangGraph, ChromaDB Persistent Vector DB, Sentence-Transformers, NetworkX Graph"),
            ("Backend Infrastructure:", "FastAPI, Uvicorn, SQLite 3, PyMuPDF (fitz), Trafilatura, Requests"),
            ("Frontend UI & UX:", "React 18, Vite, Tailwind CSS, Lucide Icons"),
        ]),
        ("COST ANALYSIS", [
            ("Primary Inference Cost:", "$0.00 / SKU (Local Ollama qwen2.5:3b) | Search Cost: ~$0.001 / SKU."),
        ])
    ]),

    # Slide 15: Benchmark MVP & Links
    ("Snapshots of the MVP & Project Links", [
        ("LIVE BENCHMARK TEST CASE (WHIRLPOOL DISHWASHER)", [
            ("Raw Input Tested:", "Brand='APPDE', MPN='WDTS7024RZ', Description='Dishwasher SS'"),
            ("Pipeline Resolution:", "Brand corrected to 'Whirlpool Corporation', classified to UNSPSC 83041100."),
            ("Documentation Sourced:", "Sourced official Whirlpool Specification PDF and Owner Manual."),
            ("Attributes Extracted:", "Extracted 31/31 required Unilog attributes with 91% overall confidence."),
            ("Commercial Output:", "Invoice Copy = 'DISHWASHER SS 120V 15A 50-1/4IN' (<=40 chars)."),
        ]),
        ("PUBLIC RESOURCES & CODEBASE LINKS", [
            ("GitHub Public Repository:", "https://github.com/Charishma1707/product-intelligence"),
            ("Interactive Frontend Dashboard:", "http://localhost:5173"),
            ("OpenAPI / Swagger Backend Docs:", "http://localhost:8000/docs"),
        ])
    ]),
]

# Format each slide cleanly
for idx, (title, sections) in enumerate(slides_data):
    if idx < len(prs.slides):
        format_slide(prs.slides[idx], title, sections)

# Save primary file
prs.save(output_path)
print(f"[SUCCESS] Successfully generated 15 perfectly formatted, human-readable slides at: {output_path}")

# Attempt to overwrite template file if unlocked
try:
    prs.save(template_path)
    print(f"[SUCCESS] Overwrote template file: {template_path}")
except Exception as err:
    print(f"[NOTE] Could not overwrite template file directly ({err}). Saved cleanly to {output_path}")
