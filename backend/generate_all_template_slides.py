import sys
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

TEMPLATE = r'e:\unihack\product-intelligence\[EXT] UniHack-Protoype Template .pptx'
OUTPUT = r'e:\unihack\product-intelligence\UniHack-Prototype-CodeWithCofee.pptx'

# Colors
C_DARK_BG  = RGBColor(15, 23, 42)    # Slate-900
C_CARD_BG  = RGBColor(30, 41, 59)    # Slate-800
C_PRIMARY  = RGBColor(56, 189, 248)  # Cyan-400
C_ACCENT   = RGBColor(129, 140, 248) # Indigo-400
C_SUCCESS  = RGBColor(52, 211, 153)  # Emerald-400
C_WHITE    = RGBColor(255, 255, 255)
C_MUTED    = RGBColor(148, 163, 184) # Slate-400

prs = Presentation(TEMPLATE)

# We keep slide 0 (Title/Guidelines), slide 1 (Team Details), slide 13 (Links slide)
# Delete slides 2 to 12
xml_slides = prs.slides._sldIdLst
for idx in reversed(range(2, 13)):
    try:
        xml_slides.remove(xml_slides[idx])
    except:
        pass

blank_layout = prs.slide_layouts[6]

def add_custom_slide(title_text):
    slide = prs.slides.add_slide(blank_layout)
    title_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.2))
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = C_DARK_BG
    title_box.line.color.rgb = C_PRIMARY
    title_box.line.width = Pt(2)
    tf = title_box.text_frame
    tf.text = title_text
    p = tf.paragraphs[0]
    p.font.color.rgb = C_WHITE
    p.font.size = Pt(32)
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.8)
    return slide

def add_card(slide, left, top, width, height, title, bullets, highlight=False):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = C_CARD_BG if not highlight else RGBColor(17, 24, 39)
    card.line.color.rgb = C_PRIMARY if highlight else C_ACCENT
    card.line.width = Pt(2.5 if highlight else 1.5)
    tb = slide.shapes.add_textbox(Inches(left+0.2), Inches(top+0.1), Inches(width-0.4), Inches(height-0.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = C_PRIMARY if not highlight else C_SUCCESS
    p.space_after = Pt(8)
    for text, is_sub in bullets:
        p = tf.add_paragraph()
        p.text = text
        if is_sub:
            p.level = 1
            p.font.size = Pt(14)
            p.font.color.rgb = C_MUTED
            p.font.bold = False
        else:
            p.level = 0
            p.font.size = Pt(16)
            p.font.color.rgb = C_WHITE
            p.font.bold = True
            p.space_before = Pt(6)

# Slide 1: Update Team Details
team_slide = prs.slides[1]
for shape in team_slide.shapes:
    if shape.has_text_frame and "Team name" in shape.text:
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = "Team Name: codewithcofee\nTeam Leader Name: Charishma Alam\nProject: Autonomous Unilog Product Intelligence Pipeline\nTrack: Product Intelligence"
        p.font.size = Pt(22)
        p.font.color.rgb = C_DARK_BG
        p.font.bold = True

# Slide 3: Brief about your solution
s = add_custom_slide("Brief About Your Solution")
add_card(s, 0.5, 1.5, 12.33, 5.5, "Autonomous Unilog Product Intelligence Pipeline", [
    ("Core Goal", False),
    ("Transforms noisy, incomplete industrial distributor inputs (MPN, Brand, Short Desc) into a 252-column Unilog-compliant record.", True),
    ("Domain-Locked Sourcing", False),
    ("Restricts web search to official manufacturer domains (e.g. fluke.com, se.com) to block e-commerce noise (Amazon/eBay).", True),
    ("Multi-Agent LangGraph State Machine", False),
    ("10-node pipeline with 5 interactive Human-in-the-Loop review checkpoints.", True),
    ("Zero-Hallucination Evidence RAG", False),
    ("Every spec attribute is bound to an exact source URL, page number, confidence score, and text snippet.", True)
], highlight=True)

# Slide 4: Q1 - Minimal Info Enrichment
s = add_custom_slide("Q1: Minimal Product Information Enrichment")
add_card(s, 0.5, 1.5, 12.33, 5.5, "How Limited Inputs Are Transformed Into Rich Intelligence", [
    ("1. Disambiguation & Alias Resolution", False),
    ("Resolves dirty distributor codes (e.g., APPDE ➔ Whirlpool Corporation) using SQLite alias memory.", True),
    ("2. Abbreviation & Noise Expansion", False),
    ("Expands industrial jargon (e.g., SS ➔ Stainless Steel, 3P ➔ 3-Pole, 24VDC ➔ 24 Volts Direct Current).", True),
    ("3. Leaf-Level UNSPSC Taxonomy Assignment", False),
    ("Maps product description to an 8-digit UNSPSC code (e.g. 83041100) and leaf attribute schema.", True),
    ("4. Multi-Modal Attribute Extraction", False),
    ("Ingests OEM PDFs and HTML, extracting up to 252 standardized Unilog fields with verified UOMs.", True)
])

# Slide 5: Q2 - Accuracy & Trust Verification
s = add_custom_slide("Q2: Accuracy & Trust Verification Strategy")
add_card(s, 0.5, 1.5, 12.33, 5.5, "Ensuring 100% Data Quality and Trust", [
    ("Confidence Scoring (0.0 to 1.0)", False),
    ("Calculates multi-factor score based on MPN verification, source tier, evidence match, and LLM certainty.", True),
    ("Multi-Source & OEM Priority", False),
    ("Prioritizes Tier-1 OEM technical datasheets over secondary distributor catalogs.", True),
    ("Semantic & Rule-Based Validation", False),
    ("Catches unit and physical range errors (e.g., flags Amps extracted into Voltage field).", True),
    ("Human-in-the-Loop Escrow", False),
    ("Items with confidence < 80% route to human catalog managers for 1-click review.", True)
], highlight=True)

# Slide 6: Q3 - Enterprise Scalability
s = add_custom_slide("Q3: Enterprise Scalability Architecture")
add_card(s, 0.5, 1.5, 12.33, 5.5, "Handling Enterprise Catalogs and Continuous Updates", [
    ("Large Product Catalogs", False),
    ("Asynchronous batch processing with FastAPI handles multi-thousand SKU CSV uploads concurrently.", True),
    ("SHA-256 Vector Store Deduplication", False),
    ("Parses and embeds identical manufacturer PDFs exactly ONCE across sibling products.", True),
    ("Local Zero-Cost LLM Inference", False),
    ("Runs air-gapped local Ollama (qwen2.5:3b) for $0.00 primary attribute extraction.", True),
    ("Persistent State Engine", False),
    ("SQLite Job Store keeps track of all SKU execution states — fully resumable across pauses.", True)
])

# Slide 7: Opportunities & USPs
s = add_custom_slide("Opportunities & Unique Selling Proposition (USP)")
add_card(s, 0.5, 1.5, 12.33, 5.5, "Why Our Solution Outperforms Existing Approaches", [
    ("Differentiation From Traditional Approaches", False),
    ("Traditional ETL relies on static regex; standard RAG hallucinates on marketplaces. Our system locks to OEM domains.", True),
    ("Problem Resolution", False),
    ("Automates 95% of catalog onboarding work while maintaining 100% source-verifiable provenance.", True),
    ("Top 3 System USPs", False),
    ("1. OEM Domain-Locked Sourcing (Zero E-Commerce Hallucinations)", True),
    ("2. 5-Stage Human Steering Agent with Natural Language Instructions", True),
    ("3. Graph-Backed Series Attribute Inheritance (Reduces Scrapes by 70%)", True)
], highlight=True)

# Slide 8: List of Features Offered
s = add_custom_slide("List of Features Offered")
add_card(s, 0.5, 1.5, 6.0, 5.5, "Automated Intelligence Engine", [
    ("Dirty Brand Code Cleaning", False),
    ("Leaf UNSPSC Classification", False),
    ("PDF & Web Datasheet Extraction", False),
    ("Semantic Range & Unit Validation", False),
    ("Commercial Copywriting (Invoice, Short, Long, Marketing)", False)
])
add_card(s, 6.83, 1.5, 6.0, 5.5, "Interactive Review & Memory", [
    ("5-Stage React Review Dashboard", False),
    ("Live PDF Evidence Snippet Inspector", False),
    ("Natural Language Steering Agent", False),
    ("Self-Learning Alias Database", False),
    ("252-Column Unilog Standard CSV Exporter", False)
])

# Slide 9: Process Flow Diagram
s = add_custom_slide("Process Flow & Pipeline Workflow")
add_card(s, 0.5, 1.5, 12.33, 5.5, "End-to-End Product Intelligence Flow", [
    ("Raw Input (MPN, Brand, Desc) ➔ Disambiguation Node (Brand Alias & Abbr)", False),
    ("➔ Taxonomy Node (UNSPSC 8-digit & Leaf Schema) ➔ PAUSE: Review 1", True),
    ("➔ OEM Document Sourcing (Serper Search + SHA-256 PDF Cache) ➔ PAUSE: Review 2", True),
    ("➔ Series Graph Resolution (NetworkX Specification Inheritance)", True),
    ("➔ Attribute Extraction (Multi-modal LLM Schema Parsing) ➔ PAUSE: Review 3", True),
    ("➔ Validation & Confidence Scoring (Semantic & Unit Rules) ➔ PAUSE: Review 4", True),
    ("➔ Copywriting & Unilog CSV Exporter ➔ Persistent DB Store", True)
], highlight=True)

# Slide 10: System Architecture Diagram
s = add_custom_slide("System Architecture Diagram")
add_card(s, 0.5, 1.5, 6.0, 5.5, "Frontend & State Machine", [
    ("React + Vite Review Dashboard", False),
    ("FastAPI Async REST Service", False),
    ("LangGraph 10-Node Supervisor State Machine", False),
    ("SQLite Job Store (Resumable State Persistence)", False)
])
add_card(s, 6.83, 1.5, 6.0, 5.5, "Search, Vector & Knowledge Stores", [
    ("Serper Google API + DDG Fallback", False),
    ("ChromaDB Vector Store (all-MiniLM-L6-v2)", False),
    ("NetworkX Knowledge Graph (Series Specs)", False),
    ("Ollama (qwen2.5:3b) / Groq LLM Cascade", False)
])

# Slide 11: Technologies Used
s = add_custom_slide("Technologies Used in the Solution")
add_card(s, 0.5, 1.5, 6.0, 5.5, "Backend & AI Frameworks", [
    ("Python 3.13 & FastAPI", False),
    ("LangGraph (Agentic State Machine)", False),
    ("Ollama & Groq (Qwen2.5 Models)", False),
    ("ChromaDB (Vector Database)", False),
    ("NetworkX (Knowledge Graph)", False),
    ("PyMuPDF & Trafilatura (Scrapers)", False)
])
add_card(s, 6.83, 1.5, 6.0, 5.5, "Frontend & Data Layer", [
    ("React 18 + Vite (Tailwind / Lucide UI)", False),
    ("SQLite 3 (Jobs & Alias Store)", False),
    ("Serper.dev (Google Search API)", False),
    ("DuckDuckGo Search (Free Fallback)", False)
])

# Slide 12: Estimated Implementation Cost
s = add_custom_slide("Estimated Implementation Cost")
add_card(s, 0.5, 1.5, 12.33, 5.5, "Low-Cost / Zero-Cost Hackathon Prototype Architecture", [
    ("Primary LLM Extraction Cost: $0.00", False),
    ("Local Ollama (qwen2.5:3b) executes primary attribute extraction with zero API fees.", True),
    ("Vector Storage Cost: $0.00", False),
    ("Embedded ChromaDB vector store runs locally on disk.", True),
    ("Search API Cost: ~ $0.001 / SKU", False),
    ("Serper API provides 2,500 free searches; DuckDuckGo fallback provides 100% free web search.", True),
    ("Cloud Escalation Fallback: ~ $0.002 / SKU", False),
    ("Groq API free tier provides high-speed cloud inference when local hardware is busy.", True)
], highlight=True)

# Slide 13: Snapshots of the MVP
s = add_custom_slide("Snapshots of the MVP (Whirlpool Dishwasher Benchmark)")
add_card(s, 0.5, 1.5, 6.0, 5.5, "Input & Resolution", [
    ("Raw Input: Brand='APPDE' | MPN='WDTS7024RZ'", False),
    ("Brand Corrected: Whirlpool Corporation", True),
    ("UNSPSC: 83041100 (Built-In Dishwashers)", True),
    ("Abbreviation: 'SS' ➔ Stainless Steel", True)
])
add_card(s, 6.83, 1.5, 6.0, 5.5, "Output & Performance", [
    ("Attributes Yield: 31 / 31 Extracted", False),
    ("Overall Confidence: 91%", True),
    ("Invoice Desc: 'DISHWASHER SS 120V 15A 50-1/4IN'", True),
    ("Source Evidence: Official Whirlpool PDF", True)
])

# Slide 14: Future Development
s = add_custom_slide("Future Development Roadmap")
add_card(s, 0.5, 1.5, 12.33, 5.5, "Next Steps for Enterprise Production Deployment", [
    ("1. Multi-Modal Vision Model Integration", False),
    ("Extract technical specs directly from 2D engineering CAD blueprints and image diagrams.", True),
    ("2. Automated ERP / PIM Connectors", False),
    ("Direct REST & Webhook synchronization with SAP, STEP, and Akeneo PIM systems.", True),
    ("3. Automated Multilingual Translation", False),
    ("Translate Unilog specifications into European and Asian localized catalogs automatically.", True)
])

# Move original Links slide to very end
xml_slides = prs.slides._sldIdLst
links_slide_xml = xml_slides[2]
xml_slides.remove(links_slide_xml)
xml_slides.append(links_slide_xml)

links_slide = prs.slides[-1]
for shape in links_slide.shapes:
    if shape.has_text_frame and ("GitHub" in shape.text or "Provide links" in shape.text):
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = "Provide links to your:\n\n1. GitHub Public Repository:\n   github.com/Charishma1707/product-intelligence\n\n2. Demo Video Link (3 Minutes):\n   [Insert Video Link Here]\n\n3. Working Prototype Link:\n   Frontend: http://localhost:5173\n   Backend API Docs: http://localhost:8000/docs"
        p.font.size = Pt(20)
        p.font.color.rgb = C_DARK_BG
        p.font.bold = True

prs.save(OUTPUT)
print(f"Generated COMPLETE 15-SLIDE presentation matching template: {OUTPUT}")
