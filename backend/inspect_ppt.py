from pptx import Presentation

prs = Presentation(r'e:\unihack\product-intelligence\[EXT] UniHack-Protoype Template .pptx')

for i, slide in enumerate(prs.slides):
    print(f"--- Slide {i+1} ---")
    for j, shape in enumerate(slide.shapes):
        if shape.has_text_frame:
            text = shape.text.replace('\n', ' ')[:100]
            print(f"Shape {j} (Type: {shape.shape_type}): {text}")
