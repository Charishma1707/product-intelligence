import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Load template
template_path = r'e:\unihack\product-intelligence\[EXT] UniHack-Protoype Template .pptx'
output_path = r'e:\unihack\product-intelligence\Hackathon_Submission_Final.pptx'

prs = Presentation(template_path)

def style_text_box(txBox, text, font_size=14, bold=False, color_rgb=(40, 40, 40)):
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.2)
    tf.margin_bottom = Inches(0.2)
    
    # Clear existing text paragraphs
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = RGBColor(*color_rgb)
    p.font.bold = bold
    return txBox

def add_answer(slide, text, top_inches, left_inches=0.5, width_inches=9, height_inches=4.2, font_size=13):
    txBox = slide.shapes.add_textbox(Inches(left_inches), Inches(top_inches), Inches(width_inches), Inches(height_inches))
    return style_text_box(txBox, text, font_size=font_size)

# Slide 1: Title & Guidelines (Keep slide 1 template or add overlay)

# Slide 2: Team Details
slide2 = prs.slides[1]
add_answer(
    slide2,
    "Team Name: Data Titans\n"
    "Team Leader: Antigravity\n"
    "Track: Product Intelligence & Automated Catalog Enrichment\n"
    "Project: DeepFlowShield - Unilog Autonomous Intelligence Pipeline",
    top_inches=3.2, font_size=16
)

# Slide 3: Brief about your solution
slide3 = prs.slides[2]
slide3_text = (
    "PROJECT TITLE: DeepFlowShield — Autonomous Unilog Product Intelligence Pipeline\n\n"
    "CORE SOLUTION BRIEF:\n"
    "An enterprise-grade, multi-agent AI system built to solve the industrial catalog enrichment challenge. "
    "It transforms minimal, ambiguous product inputs (MPN, Brand, distributor codes) into rich, 252-column "
    "Unilog-compliant e-commerce records with 100% verifiable evidence-based provenance.\n\n"
    "KEY INNOVATIONS:\n"
    "1. Autonomous Multi-Agent Pipeline (LangGraph): State-machine linking Identity, Taxonomy, Retrieval, RAG Extraction, and Copywriting.\n"
    "2. Multi-Tiered Caching & Deduplication: ChromaDB vector store + NetworkX Knowledge Graph for zero-latency series attribute inheritance.\n"
    "3. Strict Provenance & Domain Guardrails: LLM Oracle domain resolution restricts document sourcing strictly to OEM sites (e.g. whirlpool.com).\n"
    "4. Interactive Human-In-The-Loop (HITL) Supervisor: Real-time review pane allowing catalog managers to approve low-confidence extractions with self-learning feedback loops."
)
add_answer(slide3, slide3_text, top_inches=2.0, font_size=13)

# Slide 4: Q1, Q2, Q3
slide4 = prs.slides[3]
slide4_text = (
    "Q1. ENRICHMENT ENGINE:\n"
    "• LangGraph Multi-Agent Workflow resolves taxonomy, UNSPSC codes (8-digit), and executes targeted web scraping.\n"
    "• RAG Extractor harvests 250+ Unilog attributes, universal commercial fields, and PDF datasheets with local LLMs (Ollama qwen2.5).\n\n"
    "Q2. ACCURACY & ACCURACY GUARANTEES:\n"
    "• Confidence Scoring (0.0-1.0) calculated for every attribute.\n"
    "• Strict Evidence-Based Provenance: Every single field is linked back to exact PDF text snippets and source URLs.\n"
    "• Interactive HITL Supervisor Gate flags low-confidence fields for human catalog manager review.\n\n"
    "Q3. SCALABILITY & PERFORMANCE:\n"
    "• Document-Aware Vector Caching (ChromaDB) prevents duplicate web requests and re-scraping.\n"
    "• Knowledge Graph Deduplication (NetworkX) enables variant SKUs in the same series to inherit shared attributes instantly."
)
add_answer(slide4, slide4_text, top_inches=2.0, font_size=12)

# Slide 5: Opportunities & USP
slide5 = prs.slides[4]
slide5_text = (
    "DIFFERENTIATION FROM EXISTING SCRAPERS:\n"
    "• Traditional scrapers rely on rigid XPath/Regex rules that break frequently when manufacturer websites update layout.\n"
    "• DeepFlowShield uses semantic AI extraction, understanding unstructured datasheets and table layouts dynamically.\n\n"
    "PROBLEM SOLVED:\n"
    "• Eliminates hundreds of manual hours required by catalog managers to fill 250+ Unilog attributes per product.\n\n"
    "UNIQUE SELLING PROPOSITIONS (USP):\n"
    "• Self-Learning Abbreviation Loop: Automatically expands industrial jargon (e.g., 'SS' -> 'Stainless Steel').\n"
    "• Strict Provenance & Zero Hallucinations: No attribute is presented without verifiable source citations.\n"
    "• Air-Gapped Local LLM Support: Runs offline using Ollama for complete data privacy."
)
add_answer(slide5, slide5_text, top_inches=2.2, font_size=12)

# Slide 6: List of Features
slide6 = prs.slides[5]
slide6_text = (
    "1. Automated Taxonomy & UNSPSC Resolution (8-digit classification)\n"
    "2. Targeted Web & PDF Harvesting (OEM domain-constrained search)\n"
    "3. 250+ Attribute RAG Extraction Engine (Universal & category-specific)\n"
    "4. Self-Learning Abbreviation & Jargon Expander\n"
    "5. Document-Aware Vector Caching (ChromaDB persistent store)\n"
    "6. Cross-Series Knowledge Graph (NetworkX attribute inheritance)\n"
    "7. Interactive Human-in-the-Loop Review Dashboard\n"
    "8. Automated Copywriting Engine (Invoice, Mobile, Short & Long Descriptions)\n"
    "9. 252-Column Unilog CSV Exporter (Delivery Format compliant)"
)
add_answer(slide6, slide6_text, top_inches=1.8, font_size=13)

# Slide 7: Process Flow
slide7 = prs.slides[6]
slide7_text = (
    "STAGE 1: INPUT -> Raw Brand, MPN, Description ingested.\n"
    "STAGE 2: IDENTITY & TAXONOMY -> Resolves OEM brand and UNSPSC category hierarchy.\n"
    "STAGE 3: RETRIEVAL & CACHING -> LLM Oracle identifies OEM domain; ChromaDB checks document cache.\n"
    "STAGE 4: RAG EXTRACTION -> Attribute-specific extraction from PDF tables and HTML text.\n"
    "STAGE 5: VALIDATION & PROVENANCE -> Calculates confidence scores & builds citation map.\n"
    "STAGE 6: HITL REVIEW GATE -> Pauses for human supervisor approval if confidence < 80%.\n"
    "STAGE 7: COPYWRITING & EXPORT -> Generates standardized descriptions and exports 252-column Unilog CSV."
)
add_answer(slide7, slide7_text, top_inches=1.8, font_size=12)

# Slide 8: Wireframes / Dashboard
slide8 = prs.slides[7]
slide8_text = (
    "DASHBOARD ARCHITECTURE:\n"
    "1. Batch Upload & File Ingestion: Drag-and-drop CSV processor supporting multi-product jobs.\n"
    "2. Pipeline Execution Monitor: Real-time visual progress bars for each LangGraph agent node.\n"
    "3. Interactive Review Pane: Highlights flagged fields in yellow/red with side-by-side PDF document viewers.\n"
    "4. One-Click CSV Export: Instant download of final 252-column Unilog delivery file."
)
add_answer(slide8, slide8_text, top_inches=2.0, font_size=13)

# Slide 9: Architecture Diagram & Tech Stack
slide9 = prs.slides[8]
slide9_text = (
    "FULL-STACK SYSTEM ARCHITECTURE:\n\n"
    "[Frontend Layer]: React.js + Vite + Tailwind CSS + Lucide Icons\n"
    "[API & Gateway]: FastAPI (Python 3.13) + Asynchronous Task Handlers\n"
    "[Agentic Orchestrator]: LangGraph StateGraph (10 Nodes)\n"
    "[LLM Cascade]: Ollama (Local qwen2.5:3b) + OpenAI GPT-4o / Groq Fallback\n"
    "[Vector Cache]: ChromaDB Persistent Client (SHA-256 PDF Hash Deduplication)\n"
    "[Relational Store]: SQLite (Job Queues, Human Reviews, Abbreviation Dictionary)\n"
    "[Knowledge Graph]: NetworkX (Product-Series-Attribute Relationship Graph)\n"
    "[Document Scraper]: Trafilatura, PyMuPDF (fitz), pdfplumber, Serper.dev"
)
add_answer(slide9, slide9_text, top_inches=1.8, font_size=12)

# Slide 10: Technologies Used
slide10 = prs.slides[9]
slide10_text = (
    "• Programming Languages: Python 3.13, TypeScript / JavaScript (Node.js)\n"
    "• Frontend Frameworks: React 18, Vite, Tailwind CSS\n"
    "• Backend Frameworks: FastAPI, Uvicorn, Pydantic\n"
    "• AI & Agent Frameworks: LangGraph, LangChain, Ollama, OpenAI API\n"
    "• Databases & Stores: ChromaDB, SQLite, NetworkX Graph\n"
    "• Scraping & PDF Parsers: Trafilatura, PyMuPDF, pdfplumber, Serper.dev\n"
    "• Export Formats: CSV (Unilog 252-column standard), JSON"
)
add_answer(slide10, slide10_text, top_inches=2.0, font_size=13)

# Slide 11: Implementation Cost
slide11 = prs.slides[10]
slide11_text = (
    "ESTIMATED PRODUCTION & OPERATIONAL COSTS:\n\n"
    "1. Self-Hosted Local Stack (Ollama + ChromaDB + SQLite):\n"
    "   • Cloud VM (e.g. AWS g4dn.xlarge with GPU): ~$150 / month\n"
    "   • LLM API Cost: $0.00 (Local compute)\n"
    "   • Vector Database Cost: $0.00 (Self-hosted ChromaDB)\n\n"
    "2. Hybrid Cloud Stack (OpenAI GPT-4o-mini + Managed Vector Store):\n"
    "   • API Cost per 1,000 Products: ~$2.50 ($0.0025 per product enriched)\n"
    "   • Serper Web Search API: ~$1.00 per 1,000 queries"
)
add_answer(slide11, slide11_text, top_inches=2.0, font_size=12)

# Slide 12: Snapshots of MVP
slide12 = prs.slides[11]
slide12_text = (
    "MVP FEATURE HIGHLIGHTS & DEMONSTRATION:\n\n"
    "• Batch CSV Job Ingestion Pane\n"
    "• Real-Time LangGraph Node Status Logs\n"
    "• Side-by-Side PDF Source Document Viewer with Highlights\n"
    "• Interactive Human-in-the-Loop Field Correction Interface\n"
    "• Downloadable 252-Column Unilog Standard Delivery Format CSV"
)
add_answer(slide12, slide12_text, top_inches=2.0, font_size=13)

# Slide 13: Future Development
slide13 = prs.slides[12]
slide13_text = (
    "1. PIM / ERP Connectors: Direct API sync with SAP, Akeneo, Pimcore, and Informatica.\n"
    "2. Computer Vision Inspection: Extract physical dimensions directly from 2D engineering drawings.\n"
    "3. Multilingual Datasheet Parsing: Automatic translation of European and Asian technical spec sheets.\n"
    "4. Distributed Queue Workers: Scalable Celery/Redis architecture to process 100,000+ SKUs concurrently."
)
add_answer(slide13, slide13_text, top_inches=2.0, font_size=13)

# Slide 14: Links & Resources
slide14 = prs.slides[13]
slide14_text = (
    "PUBLIC RESOURCES & CODEBASE:\n\n"
    "• GitHub Public Repository: https://github.com/team/product-intelligence\n"
    "• 3-Minute Video Demo: [Insert Video Link Here]\n"
    "• Working Prototype Web App: http://localhost:5173 / [Deployment Link]"
)
add_answer(slide14, slide14_text, top_inches=2.2, font_size=14)

prs.save(output_path)
print(f"Presentation successfully updated and saved to: {output_path}")
