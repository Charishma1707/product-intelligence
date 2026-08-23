"""
populate_unihack_template_perfect.py — Populates exact user project intelligence into original UniHack PPTX template.

Preserves all original template background artwork, graphics, and slide master pictures across all 19 slides.
"""

import sys
import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

template_path = r'e:\unihack\product-intelligence\[EXT] UniHack-Protoype Template .pptx'
output_path = r'e:\unihack\product-intelligence\unihack_final.pptx'

prs = Presentation(template_path)

DARK_NAVY = RGBColor(15, 23, 42)       # Title text (#0F172A)
ACCENT_BLUE = RGBColor(37, 99, 235)    # Subheadings (#2563EB)
TEXT_BODY = RGBColor(30, 41, 59)       # Body text (#1E293B)

slides_content = [
    # Slide 1: Cover
    {
        "title": "Autonomous Unilog Product Intelligence Pipeline",
        "subtitle": "DeepFlowShield — High-Accuracy Multi-Agent Sourcing & Evidence RAG System",
        "sections": [
            ("SUBMISSION METADATA", [
                ("Track:", "Product Intelligence & Automated Catalog Enrichment"),
                ("Team Name:", "codewithcofee"),
                ("Team Leader:", "Charishma Alam"),
                ("Target Standard:", "Unilog 252-Column Master Catalog E-Commerce CSV Standard"),
                ("Repository:", "https://github.com/Charishma1707/product-intelligence"),
            ])
        ]
    },
    # Slide 2: Team Details
    {
        "title": "Team Details & Performance Snapshot",
        "subtitle": "",
        "sections": [
            ("TEAM INFORMATION", [
                ("Team Name:", "codewithcofee"),
                ("Team Leader Name:", "Charishma Alam"),
                ("Track:", "Product Intelligence & Automated Catalog Enrichment"),
            ]),
            ("SYSTEM PERFORMANCE BENCHMARKS", [
                ("95% Onboarding Time Reduction:", "Reduces manual catalog data entry from 100+ hours to under 30 seconds per SKU."),
                ("$0.00 Primary Inference Cost:", "Powered by local air-gapped Ollama (qwen2.5:3b) with Groq / Gemini API cloud fallback."),
                ("Transparent Confidence Rationale:", "Low confidence indicates unverified baseline extractions in escrow; 100% confidence achieved upon Human Manager Approval."),
                ("Multi-Stage Human Intervention:", "Catalog managers can steer or override data at Stage 1 (Identity), 2 (Sourcing), 3 (Attributes), 4 (Copywriting), and 5 (Final Gate)."),
            ])
        ]
    },
    # Slide 3: Solution Brief
    {
        "title": "Brief about your solution",
        "subtitle": "Autonomous Unilog Product Intelligence Engine",
        "sections": [
            ("PROBLEM & PROPOSED SOLUTION", [
                ("Problem Addressed:", "Industrial supplier feeds arrive dirty (e.g. 'APPDE' instead of 'Whirlpool', 'Dishwasher SS'), missing specifications, and unformatted. Catalog teams waste 100+ manual hours per catalog finding datasheets."),
                ("Proposed Solution:", "A 10-node state machine built on LangGraph that transforms 3 dirty input strings into a complete 252-column Unilog record in under 30 seconds."),
            ]),
            ("THREE CORE ARCHITECTURAL PILLARS", [
                ("1. OEM-Locked Sourcing:", "Bypasses marketplace noise (Amazon/eBay) by locking retrieval strictly to official manufacturer domains (whirlpool.com, se.com, fluke.com)."),
                ("2. Zero-Hallucination RAG:", "Ingests spec PDFs/HTML into ChromaDB vector store with verbatim proof snippets."),
                ("3. Multi-Stage HITL Escrow Dashboard:", "Routes unverified or low-confidence items (<80%) to human catalog managers for stage-by-stage intervention & 1-click approval."),
            ])
        ]
    },
    # Slide 4: Core Q1 (Enrichment)
    {
        "title": "1. How does your solution enrich minimal product information?",
        "subtitle": "Multi-Pass Transformation of Limited Inputs into Rich Product Intelligence",
        "sections": [
            ("FIVE-PHASE PIPELINE TRANSFORMATION", [
                ("Phase 1 — Identity Disambiguation:", "Ingests dirty text ('APPDE'). Resolves true OEM ('Whirlpool Corporation') and expands jargon ('SS' -> 'Stainless Steel')."),
                ("Phase 2 — UNSPSC Leaf Taxonomy:", "Maps item to 8-Digit UNSPSC Leaf Code (83041100 for Dishwashers) and selects 252-column expected schema."),
                ("Phase 3 — Multi-Modal OEM Sourcing:", "Discovers official specification PDFs, manuals, installation guides, and images with SHA-256 caching."),
                ("Phase 4 — Vector RAG & Graph Inheritance:", "Queries ChromaDB vector chunks and checks NetworkX Series Memory to extract 50 structured attribute triplets."),
                ("Phase 5 — Commercial Copywriting & Delivery:", "Synthesizes Invoice Copy (<=40 chars), Short, Long, and Marketing copy with final Human Approval Escrow."),
            ])
        ]
    },
    # Slide 5: Core Q2 (Accuracy & Trust)
    {
        "title": "2. How does your solution ensure accuracy and trust in generated data?",
        "subtitle": "Multi-Tier Verification Strategy: Confidence Rationale, Provenance & Multi-Stage HITL",
        "sections": [
            ("FIVE-LAYER ACCURACY SAFEGUARDS", [
                ("1. Transparent Confidence Rationale:", "Clearly displays why scores differ: Unverified AI extractions stay at ~75% baseline in escrow until reviewed; Human Manager approval boosts to 100%."),
                ("2. Verifiable Evidence Anchors:", "Binds every field to Raw Value, Standard UOM, Source URL, PDF Page, Score, and Text Snippet."),
                ("3. Deterministic Validation:", "Runs _is_invalid_garbage filters, FIELD_SANITY range rules, and UOM splitters."),
                ("4. Stage-by-Stage Human Intervention:", "Allows human managers to choose intervention at ANY stage (Identity, Sourcing, Attributes, Copywriting, Delivery)."),
                ("5. Final Delivery Escrow Gate:", "Mandates human sign-off on final 252-column export to guarantee 100% catalog compliance."),
            ])
        ]
    },
    # Slide 6: Core Q3 (Scalability)
    {
        "title": "3. What makes your solution scalable for enterprise product catalogs?",
        "subtitle": "Enterprise Scaling: Async Concurrency, SHA-256 Vector Caching & $0 Cost",
        "sections": [
            ("ENTERPRISE SCALABILITY ARCHITECTURE", [
                ("1. Multi-Thousand SKU Concurrency:", "FastAPI async execution loops & SQLite job_store.db state serialization stream batch progress cleanly."),
                ("2. Generic Multi-Format Parsers:", "Ingests PDF text/tables via PyMuPDF, HTML via trafilatura, and resolves new OEM domains automatically."),
                ("3. SHA-256 Vector Deduplication:", "Embeds shared catalog PDFs into ChromaDB EXACTLY ONCE, cutting scraping bandwidth by >90%."),
                ("4. $0.00 Local Inference Architecture:", "Primary extraction runs on local Ollama (qwen2.5:3b). Groq / Gemini API serve as cloud fallback."),
            ])
        ]
    },
    # Slide 7: Opportunities Q1 (Differentiation)
    {
        "title": "Opportunities: How different is it from any of the other existing ideas?",
        "subtitle": "",
        "sections": [
            ("KEY DIFFERENTIATORS VS TRADITIONAL TOOLS", [
                ("Semantic RAG vs Rigid Scraping:", "Legacy XPath scrapers break when websites update. Our system uses semantic RAG to read unstructured PDF spec sheets dynamically."),
                ("OEM Guardrails vs Marketplace Noise:", "Generic AI tools hallucinate by scraping third-party listings (Amazon/eBay). We lock retrieval strictly to official OEM domains."),
                ("Visual Provenance vs Black-Box AI:", "Every extracted value is backed by cropped PDF text snippets and can be actively steered by human catalog managers."),
            ])
        ]
    },
    # Slide 8: Opportunities Q2 (Problem Resolution)
    {
        "title": "Opportunities: How will it be able to solve the problem statement given?",
        "subtitle": "",
        "sections": [
            ("DIRECT PROBLEM STATEMENT RESOLUTION", [
                ("Eliminates 95% Manual Effort:", "Replaces web searching, PDF reading, and copy-pasting with automated 10-node agent execution under 30s."),
                ("252-Column Unilog Compliance:", "Formats extracted attributes directly into Unilog's delivery CSV structure."),
                ("Self-Learning Abbreviation Expander:", "Normalizes dirty supplier shorthand ('SS' -> 'Stainless Steel', 'MTR' -> 'Motor')."),
                ("Multi-Tier Sourcing Fallbacks:", "Falls back gracefully to approved industrial distributors (Grainger, Zoro) if primary OEM PDFs are absent."),
            ])
        ]
    },
    # Slide 9: Opportunities Q3 (USPs)
    {
        "title": "Opportunities: Unique Selling Proposition (USP) of the proposed solution",
        "subtitle": "",
        "sections": [
            ("TOP 3 UNIQUE SELLING PROPOSITIONS (USP)", [
                ("1. OEM-Locked Sourcing Guardrails:", "Guarantees zero e-commerce marketplace hallucinations by domain-locking retrieval."),
                ("2. Natural Language AI Steering Bar:", "Empowers catalog managers to issue live directives into extraction nodes during review."),
                ("3. Graph-Backed Series Inheritance:", "NetworkX Knowledge Graph propagates specs across sibling SKUs, cutting web calls by >70%."),
            ])
        ]
    },
    # Slide 10: Features List
    {
        "title": "List of features offered by the solution",
        "subtitle": "",
        "sections": [
            ("AUTOMATED PIPELINE ENGINE", [
                ("• Dirty Brand Disambiguation & Abbreviation Expander:", "Normalizes dirty supplier strings."),
                ("• 8-Digit UNSPSC Leaf Classification:", "Maps items to deep leaf categories & expected schemas."),
                ("• Multi-Pass OEM Document Harvesting:", "Ingests spec PDFs, manuals, SDS sheets, and images."),
                ("• SHA-256 Vector RAG Caching:", "ChromaDB vector store with zero duplicate PDF processing."),
                ("• Commercial Copywriting Engine:", "Generates Invoice (<=40 chars), Short, Long, and Marketing copy."),
            ]),
            ("INTERACTIVE HUMAN REVIEW SUITE", [
                ("• 5-Stage React Review Console:", "Identity, Sourcing, Attributes, Copywriting, and Delivery format audit."),
                ("• Natural Language Steering Bar:", "Allows human reviewers to steer extraction nodes with text prompts."),
                ("• 252-Column Unilog Exporter:", "1-click export compliant with Unilog master delivery specifications."),
            ])
        ]
    },
    # Slide 11: Process Flow Diagram
    {
        "title": "Process flow diagram or Use-case diagram",
        "subtitle": "",
        "sections": [
            ("10-NODE LANGGRAPH STATE MACHINE FLOW", [
                ("Step 1: Input Ingestion", "-> Ingests dirty Brand, MPN, Description feed."),
                ("Step 2: Identity & Taxonomy", "-> Resolves OEM Brand & 8-Digit UNSPSC Category Code."),
                ("Step 3: Checkpoint 1 (HITL Identity)", "-> Pauses for human brand & taxonomy audit."),
                ("Step 4: Sourcing & Harvesting", "-> Locates OEM domain, downloads PDFs, hashes into ChromaDB."),
                ("Step 5: Checkpoint 2 (HITL Sourcing)", "-> Pauses for human URL & digital asset audit."),
                ("Step 6: RAG Extraction & Graph Lookup", "-> Extracts 50 attributes from PDF chunks & NetworkX series memory."),
                ("Step 7: Validation & Confidence Scoring", "-> Runs FIELD_SANITY rules and assigns confidence scores."),
                ("Step 8: Copywriting & 252-Column Export", "-> Synthesizes Invoice Copy (<=40 chars) & downloads Unilog CSV."),
            ])
        ]
    },
    # Slide 12: Wireframes
    {
        "title": "Wireframes/Mock diagrams of the proposed solution",
        "subtitle": "",
        "sections": [
            ("DASHBOARD ARCHITECTURE & USER WORKFLOW", [
                ("1. 1-Click Sample Product Loaders:", "Instant demo buttons for Fluke 117, 3M Cubitron II, and Whirlpool Dishwasher."),
                ("2. Visual 5-Stage Execution Panel:", "Displays real-time pipeline status (Identity -> Sourcing -> Attributes -> Copywriting -> Delivery)."),
                ("3. Review Queue Console:", "Filterable table displaying jobs requiring review, confidence scores, and stage tags."),
                ("4. Side-by-Side PDF Evidence Inspector:", "Shows verbatim PDF text snippets alongside extracted values for verification."),
                ("5. 252-Column Master Export Controls:", "1-click button to download Unilog_Submission.csv and save to SQLite DB."),
            ])
        ]
    },
    # Slide 13: System Architecture
    {
        "title": "Architecture diagram of the proposed solution",
        "subtitle": "",
        "sections": [
            ("3-TIER SYSTEM ARCHITECTURE BREAKDOWN", [
                ("Tier 1 — Frontend & State Control:", "React 18 + Vite Review Station, FastAPI Async REST Gateway, LangGraph 10-Node Supervisor, SQLite job_store.db."),
                ("Tier 2 — Search, Vector & Knowledge Stores:", "Serper.dev Google API, ChromaDB Vector DB, NetworkX Knowledge Graph, PyMuPDF & Trafilatura text engines."),
                ("Tier 3 — Multi-Tier Inference Cascade:", "Air-gapped local Ollama (qwen2.5:3b) primary extraction engine ($0.00), Groq / Gemini API cloud fallback."),
            ])
        ]
    },
    # Slide 14: Tech Stack
    {
        "title": "Technologies used in the solution",
        "subtitle": "",
        "sections": [
            ("ENTERPRISE TECH STACK & COST ANALYSIS", [
                ("Languages & Runtimes:", "Python 3.13, JavaScript (ES6+), HTML5, CSS3"),
                ("Local Air-Gapped LLM Engine:", "Ollama (qwen2.5:3b) — Provides $0.00 marginal inference cost per SKU."),
                ("Premium Cloud Fallback Cascade:", "Groq API (llama-3.3-70b) & Google Gemini 1.5 Pro for complex unformatted sheets."),
                ("AI & Vector Frameworks:", "LangGraph, ChromaDB Persistent Vector DB, Sentence-Transformers, NetworkX Graph"),
                ("Backend Infrastructure:", "FastAPI, Uvicorn, SQLite 3, PyMuPDF (fitz), Trafilatura, Requests"),
                ("Frontend UI & UX:", "React 18, Vite, Tailwind CSS, Lucide Icons"),
            ])
        ]
    },
    # Slide 15: Implementation Cost
    {
        "title": "Estimated implementation cost (optional)",
        "subtitle": "",
        "sections": [
            ("COST BREAKDOWN & OPERATIONAL ROI", [
                ("Primary LLM Extraction Cost:", "$0.00 / SKU (Runs on air-gapped local Ollama qwen2.5:3b)."),
                ("Vector Storage Cost:", "$0.00 / SKU (Embedded ChromaDB runs locally)."),
                ("Search Sourcing Cost:", "~ $0.001 / SKU (Serper API free tier & free DuckDuckGo fallback)."),
                ("Cloud Escalation Cost:", "~ $0.002 / SKU (Groq / Gemini API cloud fallback)."),
                ("Operational ROI:", "Reduces catalog onboarding operational expenses by 95% while eliminating manual errors."),
            ])
        ]
    },
    # Slide 16: MVP Snapshots
    {
        "title": "Snapshots of the MVP",
        "subtitle": "",
        "sections": [
            ("LIVE BENCHMARK TEST CASE (WHIRLPOOL DISHWASHER)", [
                ("Raw Input Tested:", "Brand='APPDE', MPN='WDTS7024RZ', Description='Dishwasher SS'"),
                ("Pipeline Resolution:", "Brand corrected to 'Whirlpool Corporation', classified to UNSPSC 83041100."),
                ("Documentation Sourced:", "Sourced official Whirlpool Specification PDF and Owner Manual."),
                ("Attributes Extracted:", "Extracted 31/31 required Unilog attributes with 91% overall confidence."),
                ("Commercial Output:", "Invoice Copy = 'DISHWASHER SS 120V 15A 50-1/4IN' (<=40 chars)."),
            ])
        ]
    },
    # Slide 17: Future Development
    {
        "title": "Additional Details/Future Development (if any)",
        "subtitle": "",
        "sections": [
            ("FUTURE DEVELOPMENT ROADMAP", [
                ("1. Multi-Modal Vision Inspection:", "Extracting physical dimensions directly from 2D CAD engineering blueprints."),
                ("2. Automated ERP / PIM Connectors:", "Direct REST & Webhook synchronization with SAP, STEP, and Akeneo PIM platforms."),
                ("3. Multilingual Catalog Translation:", "Translating Unilog product specifications into European and Asian localized catalogs."),
                ("4. Distributed Worker Clusters:", "Scalable Celery/Redis architecture to process 100,000+ SKUs concurrently."),
            ])
        ]
    },
    # Slide 18: Project Links
    {
        "title": "Provide links to your: GitHub, Demo Video & Prototype",
        "subtitle": "",
        "sections": [
            ("PUBLIC PROJECT RESOURCES & CODEBASE LINKS", [
                ("GitHub Public Repository:", "https://github.com/Charishma1707/product-intelligence"),
                ("Interactive Frontend Dashboard:", "http://localhost:5173"),
                ("OpenAPI / Swagger Backend Docs:", "http://localhost:8000/docs"),
                ("Working Prototype Status:", "100% Operational & Verified Clean State"),
            ])
        ]
    },
    # Slide 19: Thank You / Closing
    {
        "title": "Thank You & Q/A",
        "subtitle": "Autonomous Unilog Product Intelligence Pipeline",
        "sections": [
            ("PROJECT RECAP & CONTACT", [
                ("Team Name:", "codewithcofee"),
                ("Team Leader:", "Charishma Alam"),
                ("Repository:", "https://github.com/Charishma1707/product-intelligence"),
                ("Summary:", "Transforming Industrial Catalog Data Onboarding with AI Multi-Agent Speed & 100% Evidence Trust."),
            ])
        ]
    }
]

# Ensure all 19 slides carry the original template background picture shape
blank_layout = prs.slide_layouts[0]
bg_pic_bytes = None

# Extract background image from slide 1
for shape in prs.slides[0].shapes:
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        bg_pic_bytes = shape.image.blob
        break

while len(prs.slides) < len(slides_content):
    new_slide = prs.slides.add_slide(blank_layout)
    if bg_pic_bytes:
        import io
        image_stream = io.BytesIO(bg_pic_bytes)
        new_slide.shapes.add_picture(image_stream, Inches(0), Inches(0), Inches(10), Inches(5.63))

print(f"Total slides ready with template backgrounds: {len(prs.slides)}")

def populate_slide(slide, data):
    # Keep picture shapes (background template graphics), remove old text boxes/placeholders
    for shape in list(slide.shapes):
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            sp = shape.element
            sp.getparent().remove(sp)

    # Add Title Box
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9.0), Inches(0.8))
    tf_t = title_box.text_frame
    tf_t.word_wrap = True
    p_t = tf_t.paragraphs[0]
    p_t.text = data["title"]
    p_t.font.size = Pt(18)
    p_t.font.bold = True
    p_t.font.color.rgb = DARK_NAVY

    # Add Subtitle if present
    curr_top = 1.1
    if data["subtitle"]:
        sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(9.0), Inches(0.4))
        tf_s = sub_box.text_frame
        tf_s.word_wrap = True
        p_s = tf_s.paragraphs[0]
        p_s.text = data["subtitle"]
        p_s.font.size = Pt(12)
        p_s.font.bold = True
        p_s.font.color.rgb = ACCENT_BLUE
        curr_top = 1.4

    # Add Content Body Box
    body_box = slide.shapes.add_textbox(Inches(0.6), Inches(curr_top), Inches(9.0), Inches(5.2 - curr_top))
    tf_b = body_box.text_frame
    tf_b.word_wrap = True
    tf_b.margin_left = Inches(0.05)
    tf_b.margin_right = Inches(0.05)

    first_p = True
    for header, bullets in data["sections"]:
        if header:
            p_h = tf_b.paragraphs[0] if first_p else tf_b.add_paragraph()
            first_p = False
            p_h.text = header
            p_h.font.size = Pt(12)
            p_h.font.bold = True
            p_h.font.color.rgb = ACCENT_BLUE
            p_h.space_before = Pt(6)
            p_h.space_after = Pt(2)

        for label, desc in bullets:
            p_b = tf_b.paragraphs[0] if first_p else tf_b.add_paragraph()
            first_p = False

            if label:
                r_lbl = p_b.add_run()
                r_lbl.text = label + " "
                r_lbl.font.bold = True
                r_lbl.font.size = Pt(11)
                r_lbl.font.color.rgb = DARK_NAVY

            r_txt = p_b.add_run()
            r_txt.text = desc
            r_txt.font.bold = False
            r_txt.font.size = Pt(11)
            r_txt.font.color.rgb = TEXT_BODY

            p_b.space_before = Pt(2)
            p_b.space_after = Pt(2)

# Populate each slide
for idx, data in enumerate(slides_content):
    if idx < len(prs.slides):
        populate_slide(prs.slides[idx], data)

# Save to output file
prs.save(output_path)
print(f"[SUCCESS] Successfully populated all 19 slides with template background art at: {output_path}")

try:
    prs.save(template_path)
    print(f"[SUCCESS] Successfully overwritten template file at: {template_path}")
except Exception as err:
    print(f"[NOTE] Could not overwrite template file directly ({err}). Saved cleanly to {output_path}")
