import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import io

def update_presentation():
    # 1. Load the original template to get the background image
    prs_template = pptx.Presentation(r'e:\unihack\product-intelligence\[EXT] UniHack-Protoype Template .pptx')
    bg_pic_bytes = None
    for shape in prs_template.slides[0].shapes:
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
            bg_pic_bytes = shape.image.blob
            break

    if not bg_pic_bytes:
        print("Error: Could not find background image in template.")
        return

    # 2. Load the Stunning Presentation
    prs = pptx.Presentation(r'e:\unihack\product-intelligence\DeepFlowShield_UniHack.pptx')
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    DARK_NAVY = RGBColor(15, 23, 42)

    # 3. Iterate through all slides to update text and inject background
    for idx, slide in enumerate(prs.slides):
        # Inject the background picture and stretch to full slide size (Fixes the squeezed image issue)
        pic = slide.shapes.add_picture(io.BytesIO(bg_pic_bytes), 0, 0, width=slide_w, height=slide_h)
        slide.shapes._spTree.insert(2, pic._element) # Move to back

        # Update text frames (font sizes, colors, and content)
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        
                        # --- TEXT REPLACEMENTS ---
                        if 'DeepFlowShield' in run.text:
                            run.text = run.text.replace('DeepFlowShield', 'Unilog Product Intelligence')

                        # Slide 2 Updates
                        if idx == 1: 
                            if '100%' in run.text and 'Evidence' not in run.text and len(run.text) < 10:
                                run.text = run.text.replace('100%', 'Transparent')
                            if 'Evidence Provenance' in run.text:
                                run.text = run.text.replace('Evidence Provenance', 'Confidence Rationale')
                            if 'Every extracted value is bound' in run.text:
                                run.text = "Low AI confidence = baseline escrow. 100% confidence requires Human Manager Approval."
                        
                        # Slide 5 Updates
                        if idx == 4: 
                            if 'Autonomous 5-Tier Confidence' in run.text:
                                run.text = "Transparent Confidence Rationale & Multi-Stage Human Intervention"
                            if '100% Human Verified' in run.text:
                                run.text = "AI extractions start low in escrow. Human Verification at any stage boosts score to 100%."
                            if '5-Stage Human-in-the-Loop' in run.text:
                                run.text = "5-Stage Human Intervention Escrow Gates"
                            if 'Pauses low-confidence items at 5 checkpoints' in run.text:
                                run.text = "Catalog managers can override data at Identity, Sourcing, Attributes, Copywriting, and Final Delivery stages."

                        # --- FONT SIZE & COLOR ADJUSTMENTS ---
                        if run.font.size is not None:
                            current_pt = run.font.size / 12700 # Convert EMUs to points
                            
                            # Make words big enough
                            if current_pt < 12:
                                run.font.size = Pt(13)
                            elif 12 <= current_pt < 16:
                                run.font.size = Pt(16)
                            elif 16 <= current_pt < 22:
                                run.font.size = Pt(22)
                            elif current_pt >= 22:
                                run.font.size = Pt(32) # Huge titles

                            # Selective Color Fixing based on Y-Position
                            # If text is in the Header (< 1.5 inches from top) or Footer (> 6.5 inches), 
                            # it sits on the white template background and MUST be Dark Navy to be visible.
                            # We leave all body text alone so it contrasts against the dark blue shape boxes.
                            shape_top_inches = shape.top / 914400 if shape.top else 0 # 1 inch = 914400 EMUs
                            if shape_top_inches < 1.6 or shape_top_inches > 6.5:
                                run.font.color.rgb = DARK_NAVY

    # 4. Save the finalized presentation
    out_path = r'e:\unihack\product-intelligence\UniHack_Final_Stunning_V3.pptx'
    prs.save(out_path)
    print(f"Successfully generated {out_path} with FULL 16:9 UniHack background and FIXED contrasts!")

if __name__ == '__main__':
    update_presentation()
