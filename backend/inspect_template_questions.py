from pptx import Presentation

pptx_path = r'e:\unihack\product-intelligence\[EXT] UniHack-Protoype Template .pptx'
prs = Presentation(pptx_path)

print(f"Total slides in template: {len(prs.slides)}")

for idx, slide in enumerate(prs.slides):
    print(f"\n--- SLIDE {idx + 1} ---")
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                print(f"[Shape text]: {text}")
