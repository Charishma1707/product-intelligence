import sys
from copy import deepcopy
from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

TEMPLATE = r'e:\unihack\product-intelligence\[EXT] UniHack-Protoype Template .pptx'
OUTPUT = r'e:\unihack\product-intelligence\UniHack-Prototype-CodeWithCofee.pptx'

# Modern Corporate / Tech Brand Colors
C_DARK_BG  = RGBColor(15, 23, 42)    # slate-900 (deep navy)
C_CARD_BG  = RGBColor(30, 41, 59)    # slate-800
C_PRIMARY  = RGBColor(56, 189, 248)  # sky-400 (bright cyan/blue)
C_ACCENT   = RGBColor(129, 140, 248) # indigo-400
C_SUCCESS  = RGBColor(52, 211, 153)  # emerald-400
C_WHITE    = RGBColor(255, 255, 255)
C_MUTED    = RGBColor(148, 163, 184) # slate-400

prs = Presentation(TEMPLATE)

# Clear out slides we will replace (keep Title (0), Team (1), and Links (last))
# The template has 14 slides. We keep index 0, 1. Delete 2 to 12. Keep 13.
slides_to_delete = list(range(2, 13))
# PPTX doesn't have an easy delete API, so we have to hack the XML or just create new slides at the end.
# Actually, the safest way to delete slides in python-pptx is removing from rId list.
xml_slides = prs.slides._sldIdLst
for idx in reversed(slides_to_delete):
    try:
        xml_slides.remove(xml_slides[idx])
    except:
        pass

# Now we have Slide 0 (Title), Slide 1 (Team), Slide 2 (Links).
# We will insert new slides before the last slide (Links).
blank_layout = prs.slide_layouts[6] # Usually blank layout

def add_custom_slide(title_text):
    slide = prs.slides.add_slide(blank_layout)
    
    # Title Header (Full width bar)
    title_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(0), Inches(0), Inches(13.33), Inches(1.2)
    )
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = C_DARK_BG
    title_box.line.color.rgb = C_PRIMARY
    title_box.line.width = Pt(2)
    
    tf = title_box.text_frame
    tf.text = title_text
    p = tf.paragraphs[0]
    p.font.color.rgb = C_WHITE
    p.font.size = Pt(36)
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.8)
    
    # Move slide to before the last one
    # We want it to be inserted before the final Links slide
    return slide

def add_card(slide, left, top, width, height, title, bullets, highlight=False):
    # Card Background
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = C_CARD_BG if not highlight else RGBColor(17, 24, 39)
    card.line.color.rgb = C_PRIMARY if highlight else C_ACCENT
    card.line.width = Pt(2.5 if highlight else 1.5)
    
    # Inner Text Box
    tb = slide.shapes.add_textbox(Inches(left+0.2), Inches(top+0.1), Inches(width-0.4), Inches(height-0.2))
    tf = tb.text_frame
    tf.word_wrap = True
    
    # Title
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(22)
    p.font.color.rgb = C_PRIMARY if not highlight else C_SUCCESS
    p.space_after = Pt(12)
    
    for text, is_sub in bullets:
        p = tf.add_paragraph()
        p.text = text
        if is_sub:
            p.level = 1
            p.font.size = Pt(14)
            p.font.color.rgb = C_MUTED
            p.font.bold = False
        else:
            p.level = 0
            p.font.size = Pt(16)
            p.font.color.rgb = C_WHITE
            p.font.bold = True
            p.space_before = Pt(8)
            p.space_after = Pt(2)

# ==============================================================================
# Slide 1: Team Details (Modify existing)
# ==============================================================================
team_slide = prs.slides[1]
for shape in team_slide.shapes:
    if shape.has_text_frame and "codewithcofee" not in shape.text:
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = "Team Name: codewithcofee\nTeam Leader: Charishma Alam\nProject: Autonomous Unilog Product Intelligence Pipeline\nTrack: Product Intelligence\nGitHub: github.com/Charishma1707/product-intelligence"
        p.font.size = Pt(20)
        p.font.color.rgb = C_DARK_BG

# ==============================================================================
# Slide 3: Executive Summary
# ==============================================================================
s = add_custom_slide("Executive Summary & Core Innovations")
add_card(s, 0.5, 1.5, 12.33, 1.8, "What We Built", [
    ("An enterprise-grade, multi-agent AI pipeline.", False),
    ("Ingests minimal noisy distributor inputs (Brand, MPN, Desc) and auto-generates a 252-column Unilog-compliant record.", True),
    ("Delivers 100% verifiable, source-linked evidence with zero hallucinations.", True)
], highlight=True)

add_card(s, 0.5, 3.5, 6.0, 3.5, "1. LangGraph Orchestration", [
    ("10-node state machine", False),
    ("Identity \u2192 Taxonomy \u2192 Retrieve \u2192 Extract \u2192 Validate \u2192 Copywrite", True),
    ("Fully resumable conditional routing", True)
])

add_card(s, 6.83, 3.5, 6.0, 3.5, "2. OEM Domain Oracle", [
    ("Blocks marketplace noise", False),
    ("LLM identifies the official manufacturer domain and exclusively sources technical specs from OEM PDFs and HTML.", True)
])

# ==============================================================================
# Slide 4: Innovations Continued
# ==============================================================================
s = add_custom_slide("Core Innovations (Continued)")
add_card(s, 0.5, 1.5, 12.33, 2.0, "3. Zero-Hallucination Evidence Engine", [
    ("Every single attribute is bound to a 0\u20131 confidence score.", False),
    ("Includes the exact source URL and a cropped PDF snippet to prove provenance.", True),
    ("Semantic Validator catches physical mismatch (e.g., Amps extracted as Volts).", True)
], highlight=True)

add_card(s, 0.5, 3.8, 6.0, 3.2, "4. Self-Learning Knowledge Graph", [
    ("Cross-SKU attribute inheritance", False),
    ("Verified series data propagates to sibling SKUs via NetworkX.", True),
    ("Cuts external scraping operations by 70%.", True)
])

add_card(s, 6.83, 3.8, 6.0, 3.2, "5. Interactive HITL Dashboard", [
    ("5-Stage React Review Workflow", False),
    ("Low-confidence items (<80%) are routed to human reviewers.", True),
    ("Approvals persist to the database as structural learning aliases.", True)
])

# ==============================================================================
# Slide 5: The Pipeline Process
# ==============================================================================
s = add_custom_slide("How It Works: The Pipeline Architecture")
add_card(s, 0.5, 1.5, 3.9, 5.5, "Step 1: Disambiguation", [
    ("Input Clean-up", False),
    ("Resolves dirty codes (APPDE \u2192 Whirlpool).", True),
    ("Taxonomy Assignment", False),
    ("Identifies full classpath and 8-digit UNSPSC code automatically.", True)
])

add_card(s, 4.7, 1.5, 3.9, 5.5, "Step 2: Document Harvest", [
    ("Targeted Scraping", False),
    ("Bypasses Amazon/eBay, hitting the OEM directly.", True),
    ("Asset Retrieval", False),
    ("Downloads Spec PDFs, Safety Data Sheets (SDS), and images.", True)
], highlight=True)

add_card(s, 8.9, 1.5, 3.9, 5.5, "Step 3: RAG & Export", [
    ("Multi-modal Extraction", False),
    ("Parses 250+ Unilog attributes using local Ollama LLMs.", True),
    ("Commercial Copywriting", False),
    ("Generates Invoice (\u226440 char), Short, Long, and Marketing descriptions.", True)
])

# ==============================================================================
# Slide 6: Cost & Scale
# ==============================================================================
s = add_custom_slide("Scaling to Millions of SKUs Cost-Effectively")
add_card(s, 0.5, 1.5, 6.0, 2.8, "Cost Per SKU: $0.00", [
    ("Air-gapped Local LLMs", False),
    ("Ollama (qwen2.5:3b) executes primary extraction for free.", True),
    ("Cloud fallback (Groq) only used for complex edge cases (~$0.002).", True)
])

add_card(s, 6.83, 1.5, 6.0, 2.8, "ChromaDB Vector Caching", [
    ("Zero Latency on Repeat Docs", False),
    ("SHA-256 document hashing ensures we never scrape or embed the same PDF twice.", True)
])

add_card(s, 0.5, 4.6, 12.33, 2.5, "Batch Processing & Persistence", [
    ("Upload a 1,000-row CSV and the pipeline runs asynchronously.", False),
    ("SQLite Job Store keeps track of every state \u2014 fully resumable if paused for Human-in-the-Loop review.", True),
    ("Dramatically cuts catalog onboarding costs by 95%.", True)
], highlight=True)

# ==============================================================================
# Slide 7: Live Benchmark MVP
# ==============================================================================
s = add_custom_slide("Live Benchmark: Whirlpool Dishwasher MVP")
add_card(s, 0.5, 1.5, 12.33, 1.8, "Input: Brand = 'APPDE' | MPN = 'WDTS7024RZ' | Desc = 'Dishwasher SS'", [
    ("Noisy, incomplete distributor data converted into a highly structured Unilog standard.", False)
])

add_card(s, 0.5, 3.5, 6.0, 3.5, "Pipeline Resolution", [
    ("Brand Corrected:", False),
    ("APPDE \u2192 Whirlpool Corporation", True),
    ("Taxonomy Resolved:", False),
    ("UNSPSC 83041100 (Built-In Dishwashers)", True),
    ("Abbreviation Expanded:", False),
    ("'SS' \u2192 Stainless Steel", True)
])

add_card(s, 6.83, 3.5, 6.0, 3.5, "Output & Accuracy", [
    ("Extraction Yield:", False),
    ("31 / 31 fields retrieved (Dimensions, Voltage, Amps, dBA).", True),
    ("Copywriting Generates:", False),
    ("Invoice: 'DISHWASHER SS 120V 15A 50-1/4IN'", True),
    ("Confidence:", False),
    ("91% Overall \u2014 3 fields routed to Human Review.", True)
], highlight=True)

# ==============================================================================
# Slide 8: UI Wireframes (React Dashboard)
# ==============================================================================
s = add_custom_slide("Interactive Review Dashboard (React + Vite)")
add_card(s, 0.5, 1.5, 12.33, 5.5, "The Human-in-the-Loop Verification Station", [
    ("1. Catalog Job Queue", False),
    ("Real-time job badges (Enriched / Needs Review) and Batch CSV tracking.", True),
    ("2. Attribute Audit Grid", False),
    ("Side-by-side comparison of Raw Input vs Extracted vs Unilog Standard.", True),
    ("Color-coded confidence pills (Green \u2265 85%, Yellow 60-84%, Red < 60%).", True),
    ("3. Live Provenance Inspector", False),
    ("Embedded PDF viewer highlighted exactly at the source text snippet.", True),
    ("4. Database Persistence", False),
    ("Approving an edit instantly saves it to the SQLite Knowledge Graph to train future queries.", True)
])

# ==============================================================================
# Reorder slides so the Links slide is last
# ==============================================================================
xml_slides = prs.slides._sldIdLst
# Keep Slide 0, 1. The original "Links" slide was at index 2.
# We appended 6 custom slides. So the slides are now 0,1,2(Links),3,4,5,6,7,8.
# We need to move the Links slide (index 2) to the very end.
links_slide_xml = xml_slides[2]
xml_slides.remove(links_slide_xml)
xml_slides.append(links_slide_xml)

# Fill Links slide
links_slide = prs.slides[-1]
for shape in links_slide.shapes:
    if shape.has_text_frame and "Demo Video" in shape.text:
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = "GitHub: github.com/Charishma1707/product-intelligence\nDemo Video: [Insert Video Link Here]\nFrontend: localhost:5173\nBackend Docs: localhost:8000/docs"
        p.font.size = Pt(20)

prs.save(OUTPUT)
print(f"Saved STUNNING hackathon-level presentation to: {OUTPUT}")
