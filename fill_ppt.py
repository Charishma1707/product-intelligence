"""
fill_ppt_v2.py — Clean, spacious, judge-ready formatting.
Keeps template question titles. Replaces content with short punchy bullets.
"""
import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

TARGET = r'e:\unihack\product-intelligence\UniHack-Prototype-CodeWithCofee.pptx'
OUTPUT = r'e:\unihack\product-intelligence\UniHack-Prototype-CodeWithCofee.pptx'

# Brand colors
C_HEADING  = RGBColor(30,  64, 175)   # deep blue
C_ACCENT   = RGBColor(99, 102, 241)   # indigo
C_BODY     = RGBColor(30,  41,  59)   # slate-900
C_MUTED    = RGBColor(71,  85, 105)   # slate-500
C_GREEN    = RGBColor(22, 163,  74)
C_WHITE    = RGBColor(255,255,255)

prs = Presentation(TARGET)

# ─────────────────────────────────────────────────────────────
# Core helper: write structured content into a shape by name
# lines = list of (text, bold, size_pt, RGBColor|None, left_indent_pt)
# ─────────────────────────────────────────────────────────────
def fill_shape(slide, shape_name, lines, wrap=True):
    for shape in slide.shapes:
        if shape.name != shape_name or not shape.has_text_frame:
            continue
        tf = shape.text_frame
        tf.word_wrap = wrap
        tf.clear()
        for idx, item in enumerate(lines):
            text, bold, size, color, space_before = item
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.space_before = Pt(space_before)
            p.space_after  = Pt(2)
            if not text:          # blank spacer line
                continue
            run = p.add_run()
            run.text = text
            run.font.bold  = bold
            run.font.size  = Pt(size)
            if color:
                run.font.color.rgb = color
        return True
    return False


def H(text, space=10):
    """Section header line."""
    return (text, True, 12, C_HEADING, space)

def B(text, space=3):
    """Bullet line."""
    return (f"  {text}", False, 10, C_BODY, space)

def SB(text, space=3):
    """Bold sub-header / key point."""
    return (f"  {text}", True, 10, C_ACCENT, space)

def GAP():
    return ("", False, 6, None, 6)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — Team Details
# ═══════════════════════════════════════════════════════════════
fill_shape(prs.slides[1], 'Google Shape;107;p26', [
    H("Team Details", 4),
    GAP(),
    SB("Team Name:   codewithcofee"),
    SB("Team Leader: Charishma Alam"),
    GAP(),
    B("Project:  Autonomous Unilog Product Intelligence Pipeline"),
    B("Track:    Product Intelligence & Automated Catalog Enrichment"),
    B("Output:   252-Column Unilog E-Commerce Delivery Standard"),
    B("GitHub:   github.com/Charishma1707/product-intelligence"),
])

# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — Brief About Your Solution
# ═══════════════════════════════════════════════════════════════
fill_shape(prs.slides[2], 'ContentBox9001', [
    H("What We Built"),
    B("An enterprise-grade, multi-agent AI pipeline that ingests minimal noisy distributor"),
    B("inputs (Brand, MPN, Description) and auto-generates a 252-column Unilog-compliant"),
    B("product record — with 100% verifiable, source-linked evidence."),
    GAP(),
    H("5 Core Innovations"),
    SB("1.  LangGraph Multi-Agent Orchestration"),
    B("10-node state machine: Identity → Taxonomy → Retrieve → Extract → Validate → Copywrite"),
    GAP(),
    SB("2.  OEM Domain Oracle"),
    B("LLM identifies the official manufacturer domain & blocks all 3rd-party reseller noise"),
    GAP(),
    SB("3.  Zero-Hallucination Evidence Engine"),
    B("Every attribute bound to a 0–1 confidence score, source URL, and raw PDF snippet"),
    GAP(),
    SB("4.  Self-Learning Knowledge System"),
    B("ChromaDB SHA-256 cache + NetworkX Knowledge Graph (70% compute reduction)"),
    GAP(),
    SB("5.  Interactive HITL Supervisor Dashboard"),
    B("5-stage React review workflow with post-approval learning persistence"),
])

# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — Q1/Q2/Q3  (3 side-by-side boxes already exist)
# ═══════════════════════════════════════════════════════════════
fill_shape(prs.slides[3], 'ContentBox9010', [
    H("How We Enrich Minimal Input"),
    GAP(),
    SB("Input: Brand + MPN + Short Description"),
    GAP(),
    SB("Step 1 — Disambiguation"),
    B("Resolves ambiguous codes (APPDE) to true OEM (Whirlpool)"),
    B("Assigns 8-digit UNSPSC + full taxonomy classpath"),
    GAP(),
    SB("Step 2 — OEM Document Harvest"),
    B("Locks to official manufacturer domain (whirlpool.com)"),
    B("Retrieves HTML specs, PDFs, SDS, product images"),
    GAP(),
    SB("Step 3 — RAG Extraction"),
    B("Local LLM extracts 250+ Unilog attributes"),
    B("Desc-Infer node expands abbreviations from description"),
    GAP(),
    SB("Step 4 — Copywriting"),
    B("Generates 6 copy channels: Invoice ≤40 chars, Mobile,"),
    B("Short, Long, Retail, Marketing + 20 feature bullets"),
])

fill_shape(prs.slides[3], 'ContentBox9011', [
    H("How We Ensure Accuracy & Trust"),
    GAP(),
    SB("OEM-Only Sourcing"),
    B("Harvester constrained to manufacturer domain only"),
    B("Zero marketplace contamination (no Amazon / eBay)"),
    GAP(),
    SB("Deterministic Confidence Scoring"),
    B("Mfr page + MPN verified = 0.98 confidence"),
    B("PDF table = 0.85  |  AI-inferred = capped at 0.60"),
    GAP(),
    SB("Anti-Hallucination Snippet Check"),
    B("Validator verifies every snippet exists in source text"),
    B("Hallucination detected → penalty –0.40 applied"),
    GAP(),
    SB("Semantic Mismatch Rules"),
    B("Catches field confusions (material in color, amps in volts)"),
    B("+ LLM plausibility audit on every extraction"),
    GAP(),
    SB("Interactive HITL"),
    B("Confidence < 80% auto-flagged for human review"),
    B("5-stage supervised workflow with learning on approval"),
])

fill_shape(prs.slides[3], 'ContentBox9012', [
    H("How We Scale to Millions of SKUs"),
    GAP(),
    SB("ChromaDB Vector Cache"),
    B("SHA-256 hashes all PDFs — repeat SKUs = 0s latency"),
    B("Zero scrape cost for documents already seen"),
    GAP(),
    SB("NetworkX Knowledge Graph"),
    B("Verified series attributes auto-inherit to sibling SKUs"),
    B("1 approval trains the graph for the entire product line"),
    GAP(),
    SB("Local LLM Cascade"),
    B("Ollama qwen2.5:3b = $0.00 per SKU extraction"),
    B("Groq cloud fallback only when local confidence < 80%"),
    GAP(),
    SB("Batch CSV Ingestion"),
    B("Upload 1,000-row distributor CSV"),
    B("Auto-processes with deduplication & in-batch caching"),
    GAP(),
    SB("SQLite Job Store"),
    B("All pipeline state persisted — fully resumable & auditable"),
])

# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — Opportunities / USP
# ═══════════════════════════════════════════════════════════════
fill_shape(prs.slides[4], 'ContentBox9002', [
    H("How Are We Different?"),
    GAP(),
    SB("Dynamic Semantic AI  vs  Rigid Scrapers"),
    B("Parses unstructured HTML/PDF without fragile XPath rules"),
    B("Survives website redesigns — no maintenance needed"),
    GAP(),
    SB("OEM Domain Lock  vs  Unfiltered Search"),
    B("Competitors scrape Amazon/eBay resellers"),
    B("We verify the true OEM domain before every harvest"),
    GAP(),
    SB("Verifiable Evidence  vs  Black-Box LLMs"),
    B("Every attribute linked to exact source text snippet"),
    B("Full clickable audit trail — zero hallucination risk"),
    GAP(),
    SB("5-Stage HITL  vs  Binary Accept/Reject"),
    B("Brand change → re-triggers retrieval"),
    B("Spec change → re-runs validation only"),
    B("Description change → re-runs copywriting only"),
])

fill_shape(prs.slides[4], 'ContentBox9003', [
    H("How We Solve the Problem Statement"),
    GAP(),
    SB("End-to-End Automation"),
    B("Noisy 6-column CSV in → 252-column Unilog CSV out"),
    B("No manual data entry required"),
    GAP(),
    SB("Handles Placeholder Garbage"),
    B("Filters 'Unbranded', 'Display Only', '--No DIB Brand--'"),
    B("Matches Unilog content guidelines exactly"),
    GAP(),
    SB("Standardizes Industrial Jargon"),
    B("Self-Learning Abbreviation Loop:"),
    B("  SST → Stainless Steel"),
    B("  BLTLN → Built-in"),
    B("  50.25in → 50-1/4 IN"),
    GAP(),
    SB("Unilog Content Rules Enforced"),
    B("Invoice ≤ 40 chars ALL CAPS, item type first"),
    B("Approved UOM abbreviations + fraction inch conversions"),
])

fill_shape(prs.slides[4], 'ContentBox9004', [
    H("Our Unique Selling Propositions"),
    GAP(),
    SB("1.  100% Evidence-Based Provenance"),
    B("Every attribute has: source URL + document type"),
    B("+ exact raw PDF text snippet — fully auditable"),
    GAP(),
    SB("2.  $0.00 Per-SKU AI Cost"),
    B("Air-gapped Ollama qwen2.5:3b + ChromaDB caching"),
    B("Eliminates API costs for 85%+ of catalog"),
    GAP(),
    SB("3.  Cross-SKU Self-Learning Knowledge Graph"),
    B("1 human approval enriches entire product series"),
    B("NetworkX auto-inherits verified attributes to all siblings"),
    GAP(),
    SB("4.  Post-Approval Learning Loop"),
    B("Every correction persists: abbreviation aliases,"),
    B("series knowledge, and attribute patterns to SQLite DB"),
])

# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — Features
# ═══════════════════════════════════════════════════════════════
fill_shape(prs.slides[5], 'ContentBox9005', [
    H("Core Pipeline Features"),
    GAP(),
    SB("1.  Autonomous Identity & Taxonomy Resolution"),
    B("Resolves ambiguous codes → true OEM + 8-digit UNSPSC"),
    GAP(),
    SB("2.  OEM Domain-Constrained Harvester"),
    B("LLM Oracle locks to official manufacturer site only"),
    B("Harvests HTML, PDFs, SDS sheets, product images"),
    GAP(),
    SB("3.  Multi-Modal RAG Attribute Extractor"),
    B("Local LLM extracts 250+ Unilog attributes"),
    B("Subcategory-ordered schema + universal commercial fields"),
    GAP(),
    SB("4.  Zero-Hallucination Confidence Engine"),
    B("Source-type-aware scoring (0.0 – 1.0)"),
    B("Snippet verification + semantic mismatch detection"),
])

fill_shape(prs.slides[5], 'ContentBox9006', [
    H("Intelligence & Output Features"),
    GAP(),
    SB("5.  Self-Learning Abbreviation Loop"),
    B("Regex + LLM expansion of industrial shorthand"),
    B("All approvals persist to desc_abbreviations table"),
    GAP(),
    SB("6.  NetworkX Knowledge Graph + ChromaDB Cache"),
    B("Series attributes auto-inherit to sibling SKUs"),
    B("SHA-256 doc hashing → 0s retrieval on repeat scrapes"),
    GAP(),
    SB("7.  5-Stage Interactive HITL Supervisor"),
    B("Identity → URLs → Attributes → Copy → Delivery Fields"),
    B("Implicit confidence boost on silent human acceptance"),
    GAP(),
    SB("8.  252-Column Unilog CSV Exporter"),
    B("6 copy channels + 20 feature bullets + dimensions"),
    B("UPC/EAN/GTIN + digital asset URLs, deduplicated by MPN"),
])

# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — Process Flow
# ═══════════════════════════════════════════════════════════════
for name, lines in [
    ('Box9100', [
        H("1.  INPUT DISTRIBUTOR DATA", 4),
        B("Brand, MPN, Short Description"),
        B("Noisy 6-column distributor CSV"),
    ]),
    ('Box9102', [
        H("2.  DISAMBIGUATION & TAXONOMY", 4),
        B("Identity Node  →  resolves true OEM (Whirlpool)"),
        B("Taxonomy Node  →  8-digit UNSPSC (83041100) + Classpath"),
    ]),
    ('Box9104', [
        H("3.  OEM HARVESTING & VECTOR CACHING", 4),
        B("LLM Oracle locks to official domain (whirlpool.com)"),
        B("ChromaDB SHA-256:  Hit = 0s / Miss = scrape + embed"),
    ]),
    ('Box9106', [
        H("4.  RAG EXTRACTION & ACCURACY ENGINE", 4),
        B("Local LLM extracts 250+ attributes from HTML + PDF"),
        B("Desc-Infer: SST → Stainless Steel from description"),
        B("Confidence Engine: 0.0–1.0 scores + PDF snippet"),
        B("Hallucination detection + semantic mismatch rules"),
    ]),
    ('Box9108', [
        H("5.  SUPERVISOR GATE & UNILOG EXPORT", 4),
        B("Confidence < 80%  →  5-Stage HITL Dashboard"),
        B("Post-approval  →  Knowledge Graph + DB learning"),
        B("Copywrite  →  6 channels  →  252-Column Unilog CSV"),
    ]),
]:
    fill_shape(prs.slides[6], name, lines)

# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — Wireframes
# ═══════════════════════════════════════════════════════════════
fill_shape(prs.slides[7], 'ContentBox9200', [
    (("REACT SUPERVISOR REVIEW DASHBOARD  —  Live MVP"), True, 13, C_HEADING, 4),
])

for name, lines in [
    ('Box9300', [
        H("Product Catalog Queue", 4),
        B("Left Panel"),
        GAP(),
        B("Real-time job status badges"),
        B("  Enriched / Review Required / Processing"),
        GAP(),
        B("Confidence score per product"),
        B("  94% High  |  72% Needs Review"),
        GAP(),
        B("Batch CSV upload with row-level tracking"),
    ]),
    ('Box9301', [
        H("Attribute Audit Grid", 4),
        B("Center Panel"),
        GAP(),
        B("Side-by-side comparison:"),
        B("  Raw Input  vs  Extracted  vs  Unilog Standard"),
        GAP(),
        B("Confidence pills per field:"),
        B("  Green >= 85%  |  Yellow 60-84%  |  Red < 60%"),
        GAP(),
        B("Inline edit + AI Agent Prompt box"),
    ]),
    ('Box9302', [
        H("Live Provenance & PDF Viewer", 4),
        B("Right Panel"),
        GAP(),
        B("Embedded PDF viewer auto-zoomed to evidence page"),
        GAP(),
        B("Highlight around exact raw text snippet used"),
        GAP(),
        B("Clickable source URL + document authority tag"),
        GAP(),
        B("'Approve & Train Graph' — saves to KG instantly"),
    ]),
]:
    fill_shape(prs.slides[7], name, lines)

# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — Architecture
# ═══════════════════════════════════════════════════════════════
for name, lines in [
    ('Box9400', [
        H("FRONTEND LAYER", 4),
        B("React 18 + Vite  —  dark glassmorphism UI"),
        B("Catalog Queue  |  HITL Review Pane  |  Provenance Inspector"),
        B("Batch Upload  |  CSV Export  |  Live Metrics Dashboard"),
    ]),
    ('Box9401', [
        H("API LAYER", 4),
        B("FastAPI (Python 3.12)  +  Pydantic v2  +  uvicorn"),
        B("Async REST: /enrich/v2  /enrich/resume  /enrich/batch"),
        B("/enrich/agent/prompt  /export/csv  /metrics  /jobs"),
    ]),
    ('Box9402', [
        H("ORCHESTRATION LAYER  —  LangGraph State Machine", 4),
        B("Identity  →  Taxonomy  →  Retrieve  →  Series"),
        B("Extract  →  Desc-Infer  →  Validate  →  Review Gate"),
        B("Copywrite  →  Finalize  (10 nodes, conditional routing)"),
    ]),
    ('Box9403', [
        H("AI & DATA LAYER", 4),
        B("Local LLM:  Ollama qwen2.5:3b  ($0.00/SKU)"),
        B("Cloud:      Groq fallback  (edge-case documents)"),
        B("ChromaDB:   PDF vector store  (SHA-256 cache)"),
        B("NetworkX:   Knowledge Graph  (series inheritance)"),
        B("SQLite:     Job store + knowledge + reviews + metrics"),
    ]),
    ('Box9404', [
        H("OUTPUT", 4),
        B("252-Column Unilog CSV Delivery File"),
        B("6 copy channels  +  20 feature bullets"),
        B("Dimensions, UPC/EAN/GTIN, digital asset URLs"),
        B("Deduplicated by MPN — ready for PIM import"),
    ]),
]:
    fill_shape(prs.slides[8], name, lines)

# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — Technologies
# ═══════════════════════════════════════════════════════════════
fill_shape(prs.slides[9], 'ContentBox9500', [
    H("AI & Orchestration"),
    SB("LangGraph"),
    B("Stateful 10-node agent workflow, conditional routing, HITL pause/resume"),
    SB("Ollama  qwen2.5:3b"),
    B("Primary RAG extractor — zero-cost, fully private, air-gapped"),
    SB("Groq  (llama-3.3-70b)"),
    B("Cloud fallback for complex edge-case documents"),
    GAP(),
    H("Data Storage & Caching"),
    SB("ChromaDB"),
    B("Vector DB for PDF embeddings + SHA-256 content-addressed cache"),
    SB("NetworkX"),
    B("In-memory Knowledge Graph for cross-SKU series attribute inheritance"),
    SB("SQLite"),
    B("Job store, HITL reviews, brand aliases, series knowledge, metrics"),
])

fill_shape(prs.slides[9], 'ContentBox9501', [
    H("Scraping & Document Parsing"),
    SB("Trafilatura + BeautifulSoup4"),
    B("Dynamic semantic HTML content extraction"),
    SB("PyMuPDF (Fitz) + pdfplumber"),
    B("Tabular extraction from technical PDF datasheets"),
    SB("DuckDuckGo Search API"),
    B("OEM domain discovery & manufacturer page lookup"),
    GAP(),
    H("Full-Stack Framework"),
    SB("Backend"),
    B("Python 3.12  |  FastAPI  |  Pydantic v2  |  uvicorn"),
    SB("Frontend"),
    B("React 18  |  Vite  |  Vanilla CSS  |  HTTP Basic Auth"),
    SB("Export"),
    B("python-pptx  |  csv  |  openpyxl  (Unilog delivery format)"),
])

# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — Cost
# ═══════════════════════════════════════════════════════════════
fill_shape(prs.slides[10], 'ContentBox9600', [
    H("Cost Per SKU"),
    SB("Standard SKU (cache hit):"),
    B("$0.00  —  Ollama local LLM + ChromaDB cached document"),
    SB("Edge-case SKU (cloud fallback):"),
    B("~$0.002  —  Groq API invoked only when local confidence < 80%"),
    SB("85%+ of catalog runs at $0.00 cost"),
    GAP(),
    H("Infrastructure (Monthly)"),
    B("GPU Cloud Host (AWS EC2 g4dn or RunPod):  ~$50–$100/month"),
    B("ChromaDB storage:  < 5 GB per 100,000 SKUs  (negligible)"),
    B("SQLite:  free, embedded, zero operational overhead"),
    GAP(),
    H("Return on Investment"),
    B("Replaces hundreds of manual data-entry hours per catalog cycle"),
    B("Cuts catalog onboarding cost by > 95%"),
    B("Knowledge Graph compounds ROI:"),
    B("  1 approved SKU trains the system for the entire product series"),
    B("  Zero marginal cost for all future sibling SKUs"),
])

# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — MVP Snapshots
# ═══════════════════════════════════════════════════════════════
fill_shape(prs.slides[11], 'ContentBox9700', [
    H("Live Benchmark  —  Whirlpool Dishwasher  (MPN: WDTS7024RZ)"),
    GAP(),
    SB("Input Given:"),
    B("Brand = APPDE  |  MPN = WDTS7024RZ  |  Desc = 'Dishwasher SS - Display Only'"),
    GAP(),
    H("Pipeline Results"),
    SB("Identity:"),
    B("APPDE  →  Whirlpool Corporation  (via LLM brand disambiguation)"),
    SB("Taxonomy:"),
    B("Appliances > Kitchen Appliances > Built-In Dishwashers  |  UNSPSC = 83041100"),
    SB("Retrieval:"),
    B("Locked to whirlpool.com  |  3 cached ChromaDB chunks  |  0s latency  |  MPN verified"),
    SB("Extraction:"),
    B("31 / 31 fields extracted  (Dimensions, Amps, dBA, Voltage, Tub Material, Series)"),
    SB("Confidence:"),
    B("Overall 91%  —  28 fields GREEN  |  3 fields YELLOW (routed to HITL)"),
    SB("Copywriting:"),
    B("Invoice = 'DISHWASHER SS 120V 15A 50-1/4IN'"),
    B("Short   = 'Whirlpool® WDTS7024RZ Built-In Dishwasher With CleanBoost™'"),
    SB("Self-Learning:"),
    B("4 abbreviation aliases saved  |  'Gold Series' boosted in Knowledge Graph"),
])

# ═══════════════════════════════════════════════════════════════
# SLIDE 13 — Future Roadmap
# ═══════════════════════════════════════════════════════════════
fill_shape(prs.slides[12], 'ContentBox9800', [
    H("Future Roadmap"),
    GAP(),
    SB("1.  Computer Vision for Engineering Schematics"),
    B("Multi-modal vision (LLaVA / GPT-4o Vision) reads dimensions"),
    B("directly from CAD drawings and exploded-view diagrams"),
    GAP(),
    SB("2.  Multi-Lingual Catalog Standardization"),
    B("Extend Abbreviation Loop to translate European & Asian"),
    B("distributor catalogs into unified English Unilog format"),
    GAP(),
    SB("3.  Scraper Hardening & Proxy Mesh"),
    B("Rotating residential proxies (BrightData) to bypass"),
    B("Cloudflare 403 blocks on restrictive manufacturer pages"),
    GAP(),
    SB("4.  Direct Unilog PIM API Connector"),
    B("One-click REST push of approved 252-column records"),
    B("into live Unilog PIM storefronts — no manual CSV uploads"),
    GAP(),
    SB("5.  Confidence Dashboard & SLA Reporting"),
    B("Executive metrics: throughput/hr, avg confidence,"),
    B("cache hit rate, HITL override rate, cost-per-SKU trends"),
])

# ═══════════════════════════════════════════════════════════════
# SLIDE 14 — Submission Links
# ═══════════════════════════════════════════════════════════════
fill_shape(prs.slides[13], 'Google Shape;179;p38', [
    H("Submission Links & Resources", 4),
    GAP(),
    SB("GitHub Public Repository:"),
    B("https://github.com/Charishma1707/product-intelligence"),
    GAP(),
    SB("Demo Video (3 Minutes):"),
    B("[Insert YouTube / Loom link here]"),
    GAP(),
    SB("Working Prototype — Frontend Dashboard:"),
    B("http://localhost:5173"),
    GAP(),
    SB("Working Prototype — FastAPI Docs:"),
    B("http://localhost:8000/docs"),
    GAP(),
    SB("Team: codewithcofee  |  Leader: Charishma Alam"),
])

# ═══════════════════════════════════════════════════════════════
prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
print("Done — clean formatted slides written.")
